"""
converter.py
------------
Kern-Logik des Prisma-Datei-Konverters.

Aufgabe: Nimmt eine Datei (PDF, DOCX, XLSX, CSV, TXT, PPTX) und wandelt sie
in das token-effizienteste Format um:
  - Fliesstext / Dokumente  -> Markdown
  - Tabellen / Daten        -> CSV

Zaehlt die Tokens der Roh-Extraktion vs. der optimierten Version, damit du
die Ersparnis siehst. Fuegt auf Wunsch XML-Tags im Anthropic-Stil hinzu.

Diese Datei hat KEINE Server-Logik - sie ist reine Konvertierung, damit du
sie auch einzeln in anderen Projekten wiederverwenden kannst.
"""

import io
import os
import sys
import csv
import re
import glob

# ---------------------------------------------------------------------------
# WINDOWS-PFADE FUER OCR-PROGRAMME AUTOMATISCH FINDEN
# ---------------------------------------------------------------------------
# Tesseract und poppler tragen sich unter Windows nicht immer in den PATH ein.
# Damit OCR trotzdem laeuft, suchen wir sie an den ueblichen Installationsorten.
# So muss der Nutzer am PATH nichts aendern.

def _frozen_base():
    """Basisordner der gebuendelten Ressourcen im EINGEFRORENEN Zustand.
    onedir/onefile: sys._MEIPASS (setzt PyInstaller), sonst der Ordner der .exe.
    Gibt None zurueck, wenn NICHT eingefroren (dann bleibt die alte Logik aktiv).
    Feste Bundle-Struktur (muss mit der .spec uebereinstimmen):
      ocr/tesseract/tesseract.exe  |  ocr/poppler/bin/  |  ocr/tessdata/"""
    if getattr(sys, "frozen", False):
        return getattr(sys, "_MEIPASS", os.path.dirname(sys.executable))
    return None


def _find_tesseract():
    """Sucht die tesseract.exe. None falls nicht da."""
    # Eingefroren: NUR den gebuendelten Ort nehmen. So faellt ein Bundling-Fehler
    # sofort auf (kein stilles Ausweichen auf eine Maschinen-Installation) und
    # auf Fremdrechnern gibt es die Maschinenpfade ohnehin nicht.
    base = _frozen_base()
    if base is not None:
        bundled = os.path.join(base, "ocr", "tesseract", "tesseract.exe")
        return bundled if os.path.isfile(bundled) else None
    # Normalbetrieb (python app.py): unveraendert die ueblichen Windows-Orte.
    candidates = [
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
        os.path.expanduser(r"~\AppData\Local\Programs\Tesseract-OCR\tesseract.exe"),
        os.path.expanduser(r"~\AppData\Local\Tesseract-OCR\tesseract.exe"),
    ]
    for c in candidates:
        if os.path.isfile(c):
            return c
    return None  # nicht gefunden -> pytesseract versucht dann den PATH


def _find_poppler_bin():
    """Sucht den poppler 'bin'-Ordner (mit pdftoppm.exe). None falls nicht da."""
    # Eingefroren: NUR den gebuendelten bin-Ordner (siehe _find_tesseract).
    base = _frozen_base()
    if base is not None:
        bundled = os.path.join(base, "ocr", "poppler", "bin")
        if os.path.isfile(os.path.join(bundled, "pdftoppm.exe")):
            return bundled
        return None
    # Normalbetrieb: unveraendert. Typische winget/manuelle Muster, '*' = Version.
    patterns = [
        r"C:\Program Files\poppler*\Library\bin",
        r"C:\Program Files\poppler*\bin",
        r"C:\poppler*\Library\bin",
        r"C:\poppler*\bin",
        os.path.expanduser(r"~\AppData\Local\Microsoft\WinGet\Packages\oschwartz10612.Poppler*\**\bin"),
        os.path.expanduser(r"~\AppData\Local\Programs\poppler*\Library\bin"),
    ]
    for pat in patterns:
        matches = glob.glob(pat, recursive=True)
        for m in matches:
            # nur akzeptieren, wenn pdftoppm.exe wirklich drin liegt
            if os.path.isfile(os.path.join(m, "pdftoppm.exe")):
                return m
    return None  # nicht gefunden -> pdf2image versucht dann den PATH


# Einmal beim Import ermitteln
_TESSERACT_PATH = _find_tesseract()
_POPPLER_PATH = _find_poppler_bin()

# Eingefroren: das gebuendelte tessdata (deu+eng+osd) aktiv setzen, damit
# lang="deu+eng" (unten) das deutsche Modell wirklich findet. Im Normalbetrieb
# NICHT setzen - dort findet Tesseract sein eigenes tessdata neben der
# Installation selbst, und wir wollen den bisherigen Betrieb nicht veraendern.
_FROZEN_BASE = _frozen_base()
if _FROZEN_BASE is not None:
    _bundled_tessdata = os.path.join(_FROZEN_BASE, "ocr", "tessdata")
    if os.path.isdir(_bundled_tessdata):
        os.environ["TESSDATA_PREFIX"] = _bundled_tessdata


# ---------------------------------------------------------------------------
# TOKEN-ZAEHLUNG
# ---------------------------------------------------------------------------
# Wir versuchen tiktoken (exakte Zaehlung, wie GPT/Claude sie intern nutzen).
# Falls tiktoken nicht verfuegbar ist (z.B. kein Internet beim ersten Start),
# fallen wir automatisch auf die bewaehrte Schaetzung ~4 Zeichen = 1 Token
# zurueck. Diese Regel wird in der Forschung fuer Englisch genutzt; fuer
# Deutsch liegt sie leicht daneben, ist aber fuer Vergleiche voellig ok.

_TIKTOKEN_ENC = None
_TIKTOKEN_TRIED = False


def _get_encoder():
    """Laedt tiktoken einmalig; gibt None zurueck falls nicht moeglich."""
    global _TIKTOKEN_ENC, _TIKTOKEN_TRIED
    if _TIKTOKEN_TRIED:
        return _TIKTOKEN_ENC
    _TIKTOKEN_TRIED = True
    try:
        import tiktoken
        _TIKTOKEN_ENC = tiktoken.get_encoding("cl100k_base")
    except Exception:
        _TIKTOKEN_ENC = None
    return _TIKTOKEN_ENC


def count_tokens(text: str) -> dict:
    """
    Zaehlt Tokens. Gibt ein dict zurueck mit:
      - count:  Anzahl Tokens
      - method: 'tiktoken' (exakt) oder 'estimate' (Schaetzung)
    """
    if not text:
        return {"count": 0, "method": "exact"}
    enc = _get_encoder()
    if enc is not None:
        return {"count": len(enc.encode(text)), "method": "tiktoken"}
    # Fallback: ~4 Zeichen pro Token
    return {"count": max(1, round(len(text) / 4)), "method": "estimate"}


# ---------------------------------------------------------------------------
# HILFSFUNKTIONEN
# ---------------------------------------------------------------------------

def _clean_text(text: str) -> str:
    """Entfernt typischen Extraktions-Muell, der nur Tokens frisst."""
    if not text:
        return ""
    # Windows-Zeilenumbrueche vereinheitlichen
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    # Mehr als 2 Leerzeilen -> genau 2 (Absatztrennung bleibt erhalten)
    text = re.sub(r"\n{3,}", "\n\n", text)
    # Trailing-Whitespace pro Zeile entfernen
    text = "\n".join(line.rstrip() for line in text.split("\n"))
    # Fuehrende/abschliessende Leerzeilen weg
    return text.strip()


def _rows_to_csv(rows) -> str:
    """Wandelt eine Liste von Zeilen (Listen) in einen CSV-String um.
    lineterminator=\\n wie alle uebrigen Exporte - der csv-Default \\r\\n
    wuerde beim Textmodus-Schreiben in app.py zu \\r\\r\\n verdoppelt."""
    buf = io.StringIO()
    writer = csv.writer(buf, lineterminator="\n")
    for row in rows:
        # None-Zellen zu leerem String machen
        writer.writerow(["" if c is None else str(c).strip() for c in row])
    return buf.getvalue().strip()


# ---------------------------------------------------------------------------
# EXTRAKTOREN PRO DATEITYP
# ---------------------------------------------------------------------------
# Jeder Extraktor gibt ein dict zurueck:
#   raw_text:     was rohe Extraktion liefern wuerde (fuer Token-Vergleich)
#   converted:    optimierter Inhalt (Markdown oder CSV)
#   target_format: 'markdown' oder 'csv'
#   note:         kurze Erklaerung fuer die Anzeige

# Ueberschriften-Muster, die NIEMALS als Kopf-/Fusszeile entfernt werden
# duerfen - auch wenn sie sich (nach Ziffern-Normalisierung) auf jeder
# Seite "wiederholen". 'Kapitel 1' bis 'Kapitel 5' sind Struktur, kein Muell.
_PROTECTED = re.compile(
    r"^\s*(kapitel|abschnitt|artikel|teil|anhang|"
    r"chapter|section|article|part|appendix)\b",
    re.IGNORECASE,
)


def _plural(n, one, many):
    """Anzeige-Zahl mit korrektem Numerus: '1 Ueberschrift' vs
    '3 Ueberschriften'. one/many sind die kompletten Wortgruppen (in der
    jeweiligen Sprache), damit auch Faelle mit veraendertem Innenteil passen
    ('1 Zeilenumbruch zu einem Absatz verbunden' vs 'K Zeilenumbrueche zu
    Absaetzen verbunden'). Sprachneutral: der Aufrufer liefert die Woerter."""
    return f"{n} {one if n == 1 else many}"


def _norm_lang(lang):
    """Validierung des Note-Sprachparameters: nur 'de'/'en', alles andere
    (None, Tippfehler, fremde Codes) faellt auf 'de' zurueck."""
    return lang if lang in ("de", "en") else "de"


# Alle nutzerseitig sichtbaren Notes (Block H). DE ist wortgleich zum
# frueheren Hardcode-Stand - keine Umformulierung nebenbei.
_NOTES = {
    "de": {
        "image_note": (" Hinweis: {n} eingebettete(s) Bild(er)/Grafik(en) "
                       "wurden nicht übernommen (nur Text ist extrahierbar)."),
        "ocr_needed": (
            "Fuer Bild-Seiten wird OCR benoetigt. Bitte installiere die Pakete:\n"
            "  pip install pdf2image pytesseract\n"
            "und das Programm Tesseract (siehe README, Abschnitt OCR)."),
        "tesseract_missing": (
            "Das Programm 'Tesseract' wurde nicht gefunden.\n"
            "Windows: Installer von https://github.com/UB-Mannheim/tesseract/wiki\n"
            "oder per winget: winget install -e --id UB-Mannheim.TesseractOCR\n"
            "Nach der Installation Tool neu starten. Details im README."),
        "poppler_missing": (
            "PDF-Seite {n} konnte nicht in ein Bild gewandelt werden ({e}).\n"
            "Windows braucht dafuer 'poppler'.\n"
            "Per winget: winget install -e --id oschwartz10612.Poppler\n"
            "Anleitung im README."),
        "bildpdf_no_text": "BILD-PDF ERKANNT (kein Text enthalten). {err}",
        "broken_skipped_prefix": ("ACHTUNG - UNVOLLSTÄNDIG: Seite(n) {pages} "
                                  "konnte(n) nicht gelesen werden und wurde(n) "
                                  "übersprungen. "),
        "rescued_all_pages": ("PDF-Seiten ließen sich nicht direkt lesen und "
                              "wurden per Texterkennung (OCR) gerettet "
                              "({n} Seiten)."),
        "bildpdf_ocr": ("BILD-PDF per OCR erkannt ({n} Seiten, Text war als "
                        "Bild eingebettet). Als Bild kostet so ein PDF ein "
                        "Vielfaches an Tokens - als Text ist es jetzt schlank "
                        "lesbar."),
        "rescued_pages": (" Seite(n) {pages} ließ(en) sich nicht direkt lesen "
                          "und wurde(n) per Texterkennung gerettet."),
        "ocr_aborted": ("ACHTUNG - UNVOLLSTÄNDIG: OCR brach ab, Seite(n) "
                        "{pages} fehlen im Ergebnis. {err} | {note}"),
        "mixed_image_fail": ("Seite(n) {pages} sind Bild-Seiten und konnten "
                             "nicht per OCR gelesen werden."),
        "mixed_broken_fail": ("Seite(n) {pages} konnte(n) nicht gelesen "
                              "werden und wurde(n) übersprungen."),
        "mixed_incomplete": ("ACHTUNG - UNVOLLSTÄNDIG: Gemischtes PDF, aber "
                             "{parts} {err} Die {n} Textseite(n) sind enthalten."),
        "mixed_ok": "Gemischtes PDF: {n} Textseite(n) direkt extrahiert.",
        "mixed_ocr_pages": (" Seite(n) {pages} per OCR erkannt (waren als "
                            "Bild eingebettet)."),
        "headers_removed": " {n} wiederkehrende Kopf-/Fusszeilen entfernt.",
        "removed_shown": " Entfernt: {shown}{more}.",
        "removed_more": " (+{n} weitere Muster)",
        "table_pdf": ("PDF enthält überwiegend Tabellen -> als CSV konvertiert "
                      "(token-effizienteste Form für Daten)."),
        "companion_kept": (" Begleittext ausserhalb der Tabellen wurde mit "
                           "übernommen."),
        "pdf_cleaned": "PDF-Fliesstext als Markdown bereinigt.",
        "pdf_plain": ("PDF-Fliesstext extrahiert und als sauberes Markdown "
                      "bereinigt."),
        "st_heading_one": "Überschrift erkannt",
        "st_heading_many": "Überschriften erkannt",
        "st_table_one": "Tabelle übernommen",
        "st_table_many": "Tabellen übernommen",
        "st_join_one": "Zeilenumbruch zu einem Absatz verbunden",
        "st_join_many": "Zeilenumbrüche zu Absätzen verbunden",
        "docx_note": ("Word-Struktur (Überschriften, Listen, Tabellen) in "
                      "natives Markdown übersetzt - in Original-Reihenfolge. "
                      "Hinweis: Fußnoten, Textboxen und Kopf-/Fußzeilen kann "
                      "die Word-Bibliothek prinzipbedingt nicht extrahieren."),
        "xlsx_note": ("Excel-Daten als CSV exportiert - das token-effizienteste "
                      "Format für Tabellen (bis zu 3x weniger als JSON/HTML)."),
        "xlsx_formula": (" Hinweis: {n} Formel-Zelle(n) ohne gespeichertes "
                         "Ergebnis - die Formel selbst wurde ausgegeben (Datei "
                         "einmal in Excel öffnen und speichern liefert die "
                         "Werte)."),
        "csv_note": ("CSV ist bereits das effizienteste Datenformat - nur "
                     "Zeilenumbrueche vereinheitlicht."),
        "txt_note": ("Textdatei bereinigt (ueberfluessige Leerzeilen und "
                     "Whitespace entfernt)."),
        "pptx_note": ("PowerPoint-Folien in Markdown umgewandelt (Text pro "
                      "Folie extrahiert)."),
        "pptx_speaker": (" Hinweis: {n} Folie(n) enthalten Sprechernotizen - "
                         "diese sind im Output mit enthalten."),
    },
    "en": {
        "image_note": (" Note: {n} embedded image(s)/graphic(s) were not "
                       "carried over (only text can be extracted)."),
        "ocr_needed": (
            "OCR is required for image pages. Please install the packages:\n"
            "  pip install pdf2image pytesseract\n"
            "and the Tesseract program (see README, OCR section)."),
        "tesseract_missing": (
            "The 'Tesseract' program was not found.\n"
            "Windows: installer from https://github.com/UB-Mannheim/tesseract/wiki\n"
            "or via winget: winget install -e --id UB-Mannheim.TesseractOCR\n"
            "Restart the tool after installation. Details in the README."),
        "poppler_missing": (
            "PDF page {n} could not be converted into an image ({e}).\n"
            "Windows needs 'poppler' for this.\n"
            "Via winget: winget install -e --id oschwartz10612.Poppler\n"
            "Instructions in the README."),
        "bildpdf_no_text": "IMAGE PDF DETECTED (no text contained). {err}",
        "broken_skipped_prefix": ("WARNING - INCOMPLETE: page(s) {pages} "
                                  "could not be read and were skipped. "),
        "rescued_all_pages": ("PDF pages could not be read directly and were "
                              "rescued via text recognition (OCR) ({n} pages)."),
        "bildpdf_ocr": ("IMAGE PDF recognized via OCR ({n} pages, text was "
                        "embedded as images). As images such a PDF costs many "
                        "times more tokens - as text it is now cheap to read."),
        "rescued_pages": (" Page(s) {pages} could not be read directly and "
                          "were rescued via text recognition."),
        "ocr_aborted": ("WARNING - INCOMPLETE: OCR aborted, page(s) {pages} "
                        "are missing from the result. {err} | {note}"),
        "mixed_image_fail": ("Page(s) {pages} are image pages and could not "
                             "be read via OCR."),
        "mixed_broken_fail": "Page(s) {pages} could not be read and were skipped.",
        "mixed_incomplete": ("WARNING - INCOMPLETE: mixed PDF, but {parts} "
                             "{err} The {n} text page(s) are included."),
        "mixed_ok": "Mixed PDF: {n} text page(s) extracted directly.",
        "mixed_ocr_pages": (" Page(s) {pages} recognized via OCR (were "
                            "embedded as images)."),
        "headers_removed": " {n} recurring header/footer lines removed.",
        "removed_shown": " Removed: {shown}{more}.",
        "removed_more": " (+{n} more patterns)",
        "table_pdf": ("PDF consists mostly of tables -> converted to CSV "
                      "(the most token-efficient form for data)."),
        "companion_kept": (" Accompanying text outside the tables was carried "
                           "over as well."),
        "pdf_cleaned": "PDF body text cleaned up as Markdown.",
        "pdf_plain": ("PDF body text extracted and cleaned up as tidy "
                      "Markdown."),
        "st_heading_one": "heading detected",
        "st_heading_many": "headings detected",
        "st_table_one": "table carried over",
        "st_table_many": "tables carried over",
        "st_join_one": "line break joined into a paragraph",
        "st_join_many": "line breaks joined into paragraphs",
        "docx_note": ("Word structure (headings, lists, tables) translated "
                      "into native Markdown - in original order. Note: "
                      "footnotes, text boxes and headers/footers cannot be "
                      "extracted by the Word library."),
        "xlsx_note": ("Excel data exported as CSV - the most token-efficient "
                      "format for tables (up to 3x less than JSON/HTML)."),
        "xlsx_formula": (" Note: {n} formula cell(s) without a cached result "
                         "- the formula itself was written out (opening and "
                         "saving the file once in Excel provides the values)."),
        "csv_note": ("CSV is already the most efficient data format - only "
                     "line endings were unified."),
        "txt_note": ("Text file cleaned up (superfluous blank lines and "
                     "whitespace removed)."),
        "pptx_note": ("PowerPoint slides converted to Markdown (text "
                      "extracted per slide)."),
        "pptx_speaker": (" Note: {n} slide(s) contain speaker notes - these "
                         "are included in the output."),
    },
}


def _nt(lang, key, **kw):
    """Note-Text in der gewuenschten Sprache, formatiert."""
    return _NOTES[_norm_lang(lang)][key].format(**kw)


# Alle nutzerseitig sichtbaren FEHLERTEXTE (Block I) - gleiche Mechanik wie
# _NOTES, aber eigener Katalog (Fehler != Notes; app.py nutzt ihn mit).
# DE ist wortgleich zum frueheren Hardcode-Stand; beschlossene Ausnahmen:
# 'unsupported_type' (auf die app.py-Fassung vereinheitlicht) und
# 'batch_too_many' (an den Client-Wortlaut T().batchTooMany angeglichen).
_ERRORS = {
    "de": {
        "unsupported_type": ("Dateityp '{ext}' nicht unterstuetzt. "
                             "Moeglich: {exts}"),
        "read_failed": "Fehler beim Lesen der Datei: {e}",
        "pdf_unreadable": ("Keine der {n} Seite(n) dieser PDF konnte gelesen "
                           "werden - auch nicht per Texterkennung. Die Datei "
                           "ist vermutlich beschädigt."),
        "filename_missing": "Dateiname fehlt.",
        "no_file": "Keine Datei empfangen.",
        "no_files": "Keine Dateien empfangen.",
        "batch_too_many": ("Maximal {n} Dateien gleichzeitig — bitte in "
                           "kleineren Gruppen konvertieren."),
        "batch_item_failed": "Fehler beim Verarbeiten der Datei: {e}",
        "unknown_error": "Unbekannter Fehler.",
        "zip_no_files": "Keine Dateien zum Herunterladen angegeben.",
        "zip_none_found": ("Keine gueltigen Dateien fuer den Download "
                           "gefunden (evtl. Server neu gestartet)."),
        "open_windows_only": ("Ordner-Oeffnen wird nur unter Windows "
                              "unterstuetzt."),
        "open_failed": "Ordner konnte nicht geoeffnet werden: {e}",
        "download_not_found": ("Datei nicht gefunden "
                               "(evtl. Server neu gestartet)."),
        "too_large": ("Datei bzw. Dateien zusammen zu gross "
                      "(Limit: {mb} MB pro Vorgang). Erhoehe MAX_MB in "
                      "app.py, falls du groessere Dateien brauchst."),
    },
    "en": {
        "unsupported_type": ("File type '{ext}' not supported. "
                             "Possible: {exts}"),
        "read_failed": "Error reading the file: {e}",
        "pdf_unreadable": ("None of the {n} page(s) of this PDF could be "
                           "read - not even via text recognition. The file "
                           "is probably damaged."),
        "filename_missing": "Filename missing.",
        "no_file": "No file received.",
        "no_files": "No files received.",
        "batch_too_many": ("At most {n} files at once — please convert in "
                           "smaller groups."),
        "batch_item_failed": "Error processing the file: {e}",
        "unknown_error": "Unknown error.",
        "zip_no_files": "No files specified for download.",
        "zip_none_found": ("No valid files found for download "
                           "(the server may have been restarted)."),
        "open_windows_only": ("Opening the folder is only supported "
                              "on Windows."),
        "open_failed": "The folder could not be opened: {e}",
        "download_not_found": ("File not found "
                               "(the server may have been restarted)."),
        "too_large": ("File or files together too large "
                      "(limit: {mb} MB per operation). Increase MAX_MB in "
                      "app.py if you need larger files."),
    },
}


def _et(lang, key, **kw):
    """Fehlertext in der gewuenschten Sprache, formatiert."""
    return _ERRORS[_norm_lang(lang)][key].format(**kw)


def _removed_note(removed_lines, max_show=4, lang="de"):
    """Formatiert die entfernten Kopf-/Fusszeilen fuer den Hinweistext,
    damit der Nutzer Fehlgriffe sofort sehen kann (Transparenz)."""
    if not removed_lines:
        return ""
    shown = ", ".join(f"'{l}'" for l in removed_lines[:max_show])
    more = ""
    if len(removed_lines) > max_show:
        more = _nt(lang, "removed_more", n=len(removed_lines) - max_show)
    return _nt(lang, "removed_shown", shown=shown, more=more)


def _structure_note(n_headings, n_tables, n_joins, lang="de"):
    """Beziffert die geleistete Strukturarbeit fuer die Anzeige-Note (Block C).
    Nur Zaehler > 0 werden genannt, damit die Note bei einfachen PDFs nicht
    unnoetig laenger wird (0/0/0 -> leerer String)."""
    parts = []
    if n_headings > 0:
        parts.append(_plural(n_headings, _nt(lang, "st_heading_one"),
                             _nt(lang, "st_heading_many")))
    if n_tables > 0:
        parts.append(_plural(n_tables, _nt(lang, "st_table_one"),
                             _nt(lang, "st_table_many")))
    if n_joins > 0:
        parts.append(_plural(n_joins, _nt(lang, "st_join_one"),
                             _nt(lang, "st_join_many")))
    return ", ".join(parts) + "." if parts else ""


def _strip_repeating_headers_footers(pages_lines, key=None):
    """
    Entfernt Zeilen, die auf vielen Seiten identisch wiederkehren
    (typische Kopf-/Fusszeilen wie 'Seite X', Firmennamen, Copyright).
    Das ist echter, entfernbarer Token-Muell in mehrseitigen PDFs.

    Sicherheitsregeln gegen Inhaltsverlust:
      - Nur die ersten/letzten EDGE Zeilen einer Seite sind Kandidaten
        (Kopf-/Fusszeilen stehen am Seitenrand, Inhalt in der Mitte).
      - Ueberschriften-Muster (_PROTECTED, z.B. 'Kapitel 3') sind tabu.
      - Ziffern werden nur bei kurzen Zeilen (<=40 Zeichen) normalisiert,
        damit 'Seite 1'/'Seite 2' als ein Muster gelten - laengere Zeilen
        muessen exakt uebereinstimmen.

    pages_lines: Liste von Seiten, jede Seite eine Liste von Zeilen-Items.
    key: optionaler Accessor Item -> Zeilentext. Ohne key ist das Item
         selbst der Text (bisheriges Verhalten). Mit key koennen Zeilen-
         Records (Text + Geometrie, Block C) gefiltert werden, ohne dass
         die Wiederholungs-Logik davon etwas merkt.
    Rueckgabe: (bereinigte_seiten, anzahl_entfernter_zeilen, entfernte_muster)
    - bereinigte_seiten hat dieselbe Form wie die Eingabe.
    """
    from collections import Counter

    if key is None:
        key = lambda line: line

    n_pages = len(pages_lines)
    if n_pages < 2:
        # Bei 1 Seite gibt es keine "Wiederholung" -> nichts entfernen
        return pages_lines, 0, []

    EDGE = 2            # nur die aeussersten 2 Zeilen oben/unten pruefen
    MAX_HEADER_LEN = 80  # laengere Zeilen sind sicher Inhalt

    def normalize(text):
        s = text.strip()
        if len(s) <= 40:
            return re.sub(r"\d+", "#", s)
        return s

    def is_candidate(idx, page_len, text):
        s = text.strip()
        if not s or len(s) > MAX_HEADER_LEN:
            return False
        if _PROTECTED.match(s):
            return False
        return idx < EDGE or idx >= page_len - EDGE

    counter = Counter()
    for page in pages_lines:
        # pro Seite nur eindeutige Muster zaehlen (sonst zaehlt Wiederholung
        # innerhalb einer Seite mit)
        seen = set()
        for idx, line in enumerate(page):
            if not is_candidate(idx, len(page), key(line)):
                continue
            norm = normalize(key(line))
            if norm and norm not in seen:
                seen.add(norm)
                counter[norm] += 1

    # Ein Muster gilt nur dann als Kopf/Fusszeile, wenn es auf FAST ALLEN
    # Seiten vorkommt (>=80%). Das schuetzt vor Fehlalarmen.
    threshold = max(2, n_pages * 0.8)
    junk_patterns = {norm for norm, cnt in counter.items() if cnt >= threshold}

    removed = 0
    removed_lines = []
    cleaned_pages = []
    for page in pages_lines:
        kept = []
        for idx, line in enumerate(page):
            text = key(line)
            if is_candidate(idx, len(page), text) and normalize(text) in junk_patterns:
                removed += 1
                stripped = text.strip()
                if stripped not in removed_lines:
                    removed_lines.append(stripped)
                continue
            kept.append(line)
        cleaned_pages.append(kept)

    return cleaned_pages, removed, removed_lines


# ---------------------------------------------------------------------------
# PDF-STRUKTUR (Block C): ZEILEN-RECORDS UND UEBERSCHRIFTEN
# ---------------------------------------------------------------------------
# Grundidee: pdfplumber liefert ueber extract_text_lines() dieselben Zeilen
# wie extract_text(), aber MIT Geometrie (x0/x1/top) und Zeichendaten
# (Schriftgroesse, Fontname). Die Strukturlogik arbeitet auf diesen
# annotierten Zeilen; raw_text (Messbasis "vorher") bleibt unveraendert der
# nackte extract_text()-Join.

# Ueberschriften-Erkennung: KONSERVATIV. Ein falsches Heading ist schlimmer
# als keins - deshalb muessen ALLE Signale gleichzeitig zustimmen.
_HEADING_SIZE_FACTOR = 1.2   # Zeilengroesse >= 1.2x Fliesstext-Median
_HEADING_BOLD_MIN = 0.6      # Mehrheit der Zeichen der ZEILE fett
_HEADING_WIDTH_MAX = 0.85    # Zeile deutlich schmaler als der Satzspiegel
_HEADING_MAX_CHARS = 100     # laengere Zeilen sind sicher Fliesstext
_HEADING_SIZE_CLUSTER = 0.5  # pt-Toleranz: Groessen in einem Cluster = eine Ebene


def _line_records(page, expected_text):
    """Zeilen einer Seite als Records mit Geometrie- und Font-Daten.

    Sicherheitsanker des ganzen Ansatzes: Ergeben die Zeilen NICHT exakt
    den extract_text()-Text, geben wir None zurueck - die Seite laeuft dann
    wie bisher als nackte Textzeilen (kein Inhaltsrisiko, nur keine
    Struktur-Extras). Suite 13 beweist die Gleichheit fuer den Testkorpus.
    """
    try:
        lines = page.extract_text_lines(return_chars=True)
    except Exception:
        return None
    if "\n".join(l["text"] for l in lines) != expected_text:
        return None
    recs = []
    for l in lines:
        chars = l.get("chars") or []
        sizes = sorted(round(c.get("size", 0.0), 1) for c in chars)
        visible = [c for c in chars if (c.get("text") or "").strip()]
        n_bold = sum(1 for c in visible
                     if "bold" in (c.get("fontname") or "").lower())
        recs.append({
            "text": l["text"],
            "x0": l["x0"], "x1": l["x1"],
            "top": l["top"], "bottom": l["bottom"],
            "size": sizes[len(sizes) // 2] if sizes else None,
            "bold_ratio": (n_bold / len(visible)) if visible else 0.0,
            "n_chars": len(visible),
        })
    return recs


def _plain_records(text):
    """Zeilen ohne Geometrie (OCR-Seiten, Fallback) im selben Record-Format."""
    return [{"text": s} for s in text.split("\n")]


# Toleranz beim Zuordnen von Zeilen zu Tabellen-BBoxes (Rundungsartefakte
# an den Gitterlinien, vgl. den Newmont-Ueberhang aus Block H)
_TABLE_BBOX_TOL = 2.0


def _mark_table_zones(pages_items, pages_tables):
    """Markiert Zeilen, die in einer Tabellen-BBox liegen (rec['table']=k).

    Zonen-Semantik wie pdfplumber.outside_bbox (das der CSV-Zweig seit
    Block 1 nutzt): ein Objekt gehoert zur Zone, wenn es VOLLSTAENDIG in
    der BBox liegt. Eine Zeile liegt genau dann vollstaendig drin, wenn
    ihre eigene BBox drin liegt - das wenden wir auf die konsistenz-
    gesicherten Zeilen-Records an, statt einen zweiten Extraktionslauf
    ueber outside_bbox zu starten (der wuerde den extract_text-Anker aus
    C2 brechen). Ragt eine Zeile seitlich ueber die Tabelle hinaus
    (Mehrspalten-Faelle, ausserhalb des Scopes), bleibt sie konservativ
    Fliesstext - Textverlust ist das schlimmere Uebel als eine Dublette.
    """
    for page_recs_, tables in zip(pages_items, pages_tables):
        if not tables:
            continue
        for rec in page_recs_:
            if rec.get("top") is None:
                continue
            for k, t in enumerate(tables):
                bx0, btop, bx1, bbottom = t["bbox"]
                if (rec["x0"] >= bx0 - _TABLE_BBOX_TOL
                        and rec["x1"] <= bx1 + _TABLE_BBOX_TOL
                        and rec["top"] >= btop - _TABLE_BBOX_TOL
                        and rec["bottom"] <= bbottom + _TABLE_BBOX_TOL):
                    rec["table"] = k
                    break


def _pipe_table(rows):
    """Tabelle als Markdown-Pipe-Tabelle. Zell-Robustheit:
    '|' -> '\\|' (sonst zerbricht die Tabelle), Zeilenumbruch in der
    Zelle -> Leerzeichen (Pipe-Zellen sind einzeilig), None/leer ->
    leere Zelle. Alle Zeilen werden auf die Header-Breite normalisiert."""
    def cell(c):
        s = "" if c is None else str(c)
        return s.replace("\n", " ").replace("|", "\\|").strip()

    norm = [[cell(c) for c in row] for row in rows]
    width = len(norm[0])
    lines = ["| " + " | ".join(norm[0]) + " |",
             "| " + " | ".join(["---"] * width) + " |"]
    for row in norm[1:]:
        row = (row + [""] * width)[:width]
        lines.append("| " + " | ".join(row) + " |")
    return lines


def _mark_headings(pages_items):
    """Markiert konservativ erkannte Ueberschriften (rec['heading'] = 1..3).

    Nur Zeilen mit Geometrie kommen infrage (OCR-Seiten nie). Eine Zeile
    wird NUR dann Heading, wenn ALLE Signale zustimmen:
      - Schriftgroesse >= _HEADING_SIZE_FACTOR x Fliesstext-Median des
        Dokuments (zeichengewichteter Median - der Fliesstext dominiert)
      - Mehrheit der Zeichen der Zeile fett (Zeilenkontext! Ein fettes
        WORT im Fliesstext kippt die Zeile nicht)
      - Zeile deutlich schmaler als der Satzspiegel (Headings fuellen
        keine Blocksatz-Zeile) und <= _HEADING_MAX_CHARS Zeichen
      - endet nicht auf Satzzeichen, enthaelt mindestens einen Buchstaben
    Ebenen: Groessencluster absteigend -> # / ## / ###. Kandidaten jenseits
    des dritten Clusters bleiben konservativ Fliesstext.
    Rueckgabe: Anzahl markierter Ueberschriften.
    """
    # Tabellenzonen (C3) sind fuer Headings tabu und verfaelschen auch den
    # Fliesstext-Median nicht (Zellen sind oft kleiner gesetzt).
    geo = [r for page in pages_items for r in page
           if r.get("size") and r.get("n_chars") and r.get("table") is None]
    if not geo:
        return 0

    # Zeichengewichteter Median der Zeilengroessen = Fliesstext-Groesse
    weighted = sorted((r["size"], r["n_chars"]) for r in geo)
    total = sum(w for _, w in weighted)
    acc = 0
    body_size = weighted[-1][0]
    for s, w in weighted:
        acc += w
        if acc * 2 >= total:
            body_size = s
            break

    frame_w = max(r["x1"] for r in geo) - min(r["x0"] for r in geo)
    if frame_w <= 0 or body_size <= 0:
        return 0

    candidates = []
    for page in pages_items:
        for r in page:
            if not r.get("size") or not r.get("n_chars"):
                continue
            if r.get("table") is not None:
                continue  # Tabellenzonen-Veto (C3)
            text = r["text"].strip()
            if not text or len(text) > _HEADING_MAX_CHARS:
                continue
            if not any(ch.isalpha() for ch in text):
                continue
            if text[-1] in ".!?;:,":
                continue
            if r["size"] < body_size * _HEADING_SIZE_FACTOR:
                continue
            if r["bold_ratio"] < _HEADING_BOLD_MIN:
                continue
            if (r["x1"] - r["x0"]) > frame_w * _HEADING_WIDTH_MAX:
                continue
            candidates.append(r)
    if not candidates:
        return 0

    # Groessencluster (0.5-pt-Toleranz) -> maximal 3 Ebenen, groesste = #
    sizes = sorted({r["size"] for r in candidates}, reverse=True)
    clusters = [[sizes[0]]]
    for s in sizes[1:]:
        if clusters[-1][-1] - s <= _HEADING_SIZE_CLUSTER:
            clusters[-1].append(s)
        else:
            clusters.append([s])
    level_of = {}
    for level, cluster in enumerate(clusters[:3], start=1):
        for s in cluster:
            level_of[s] = level

    n = 0
    for r in candidates:
        level = level_of.get(r["size"])
        if level:
            r["heading"] = level
            n += 1
    return n


# ---------------------------------------------------------------------------
# C4: ZEILENVERBINDUNG UND SILBENTRENNUNG
# ---------------------------------------------------------------------------
# Kernregel: Ein Umbruch wird NUR entfernt, wenn Geometrie UND Text zustimmen.
# Im Zweifel bleibt der Umbruch erhalten - stiller Strukturverlust ist das
# schlimmere Uebel als ein behaltener Umbruch (dieselbe konservative
# Asymmetrie wie bei der Heading-Erkennung in C2).
#
# GEOMETRIE (notwendig): die OBERE Zeile muss "voll" sein, also bis an den
# rechten Textrand der Seite reichen - dann hat sie der Setzer umbrochen, nicht
# der Autor beendet. Eine kurze Zeile erreicht den Rand nie und zieht deshalb
# nie einen Nachfolger; damit IST der Voll-Test zugleich die Kurzzeilen-
# Sicherung (Adressbloecke, abgesetzte Kurzzeilen bleiben stehen), ein
# separates Kurzzeilen-Veto ist unnoetig. Auf derselben Seite muss die
# naechste Zeile ausserdem im selben Absatz-Block liegen (vertikaler Abstand
# nicht groesser als das Leading mal _JOIN_BLOCK_FACTOR).
#
# TEXT: zwei Signale mit BEWUSST unterschiedlichem Gewicht -
#   STARK, IMMER Pflicht (in JEDEM Regime): die obere Zeile endet NICHT auf
#     einem Satzschlusszeichen (. ! ? : ;). Faengt die Falle "geometrisch voll,
#     aber inhaltlich abgeschlossen" (ein Satz endet genau am Justierungsrand).
#     Dieses Signal wird NIRGENDS umgangen - auch nicht im Blocksatz (Test
#     test_c4_blocksatz_period_not_bypassed sichert das negativ ab).
#   SCHWACH, nur im Flattersatz Pflicht: die Folgezeile beginnt klein. Im
#     Deutschen ist das ein SCHWACHES Signal, weil jedes Substantiv gross
#     geschrieben wird und Zeilen sehr oft vor einem Substantiv umbrechen -
#     als harte Pflicht wuerde es einen Grossteil korrekter Umbrueche brechen
#     (z.B. "...gerne zur" + "Verfuegung"). Deshalb nur verlangt, wenn der
#     rechte Rand unscharf ist (Flattersatz). Bei echtem Blocksatz (die obere
#     Zeile sitzt auf einem von >=2 Zeilen geteilten Justierungsrand) beweist
#     die Geometrie den Umbruch bereits, und das schwache Signal wird
#     fallengelassen. Hintergrund: Vault-Notiz "Blocksatz-Bypass fuer deutsche
#     Substantive bei PDF-Zeilenverbindung".
#
# Alle Toleranzen sind aus den Fixture-Messungen (Suite 13) hergeleitet.

_JOIN_RIGHT_EDGE_TOL = 2.0   # pt: Zeile gilt als "voll", wenn ihr x1 so nah am
                             # rechten Textrand liegt. Gemessen: alle ziehenden
                             # Zeilen sitzen bei gap 0.00, die naechste NICHT
                             # ziehende (Fixture b, 28%-Zeile) erst bei gap 276 -
                             # dazwischen ist alles leer. Bewusst = _TABLE_BBOX_TOL.
_JOIN_BLOCK_FACTOR = 1.35    # top-Delta > Leading*Faktor => neuer Absatz-Block.
                             # Gemessen nur zwei Deltas: intra 14.0, inter 24.0;
                             # Schwelle 18.9 bei Leading 14 liegt fast symmetrisch
                             # dazwischen (26% unter intra, 27% ueber inter).
_JUSTIFY_SHARE_TOL = 0.5     # pt: teilen sich >=2 Body-Zeilen den rechten Rand
                             # so eng, gilt er als echte Blocksatz-Justierungskante.
_JOIN_TERMINAL_PUNCT = ".!?:;"


def _median(xs):
    s = sorted(xs)
    n = len(s)
    if n == 0:
        return None
    return s[n // 2] if n % 2 else (s[n // 2 - 1] + s[n // 2]) / 2


# Listen-Marker am Zeilenanfang: solche Zeilen verbinden nie (weder ziehen sie
# noch werden sie gezogen) - Aufzaehlungen bleiben zeilenweise.
_LIST_MARKER_RE = re.compile(
    r"^\s*([•‣●◦⁃∙*\-–—]\s|\d+[.)]\s)")


def _is_list_marker(text):
    return bool(_LIST_MARKER_RE.match(text or ""))


def _page_geometry(page):
    """Per-Seite-Geometrie fuer die Zeilenverbindung (C4).
    right_edge:      rechter Textrand (max x1 der Body-Zeilen).
    justify_margin:  right_edge, WENN ihn >=2 Body-Zeilen innerhalb
                     _JUSTIFY_SHARE_TOL teilen (echter Blocksatzrand), sonst None.
    leading:         Median der vertikalen Abstaende UNMITTELBAR benachbarter
                     Body-Zeilen (Heading/Tabelle unterbricht die Nachbarschaft).
    deltas:          die gemessenen Abstaende (fuer den dokumentweiten Fallback).
    Body = Zeile mit Geometrie, kein Heading, keine Tabellenzone. Seiten ohne
    Geometrie (OCR/Fallback) liefern ueberall None -> keine Verbindung.
    """
    body = [r for r in page
            if r.get("x1") is not None and r.get("heading") is None
            and r.get("table") is None]
    right_edge = max((r["x1"] for r in body), default=None)
    justify_margin = None
    if right_edge is not None:
        n_at = sum(1 for r in body if right_edge - r["x1"] <= _JUSTIFY_SHARE_TOL)
        if n_at >= 2:
            justify_margin = right_edge
    deltas = []
    prev_body = None
    for r in page:
        if (r.get("x1") is None or r.get("heading") is not None
                or r.get("table") is not None):
            prev_body = None
            continue
        if (prev_body is not None and r.get("top") is not None
                and prev_body.get("top") is not None):
            d = r["top"] - prev_body["top"]
            if d > 0:
                deltas.append(d)
        prev_body = r
    return {"right_edge": right_edge, "justify_margin": justify_margin,
            "leading": _median(deltas), "deltas": deltas}


def _line_full(rec, geo):
    """Zeile reicht bis an den rechten Textrand (vom Setzer umbrochen)?"""
    edge = geo["right_edge"]
    return (edge is not None and rec.get("x1") is not None
            and edge - rec["x1"] <= _JOIN_RIGHT_EDGE_TOL)


def _line_at_margin(rec, geo):
    """Zeile sitzt auf dem geteilten Blocksatz-Justierungsrand (falls vorhanden)
    - dann ist der Umbruch geometrisch bewiesen (Blocksatz-Bypass)."""
    jm = geo["justify_margin"]
    return (jm is not None and rec.get("x1") is not None
            and abs(rec["x1"] - jm) <= _JUSTIFY_SHARE_TOL)


def _merge_join(acc, add):
    """Verbindet zwei Zeilen. Silbentrennung am Zeilenende:
    endet die obere Zeile auf '-', wird der Bindestrich ENTFERNT, wenn davor ein
    Buchstabe steht UND die Folgezeile klein beginnt (weiche Trennung:
    'Silben-'+'trennung' -> 'Silbentrennung'); sonst BLEIBT er (echtes
    Bindestrich-Wort: 'E-'+'Mail' -> 'E-Mail'). Ohne Bindestrich: Leerzeichen."""
    left = acc.rstrip()
    right = add.lstrip()
    if left.endswith("-"):
        core = left[:-1].rstrip()
        if core and core[-1].isalpha() and right[:1].islower():
            return core + right
        return left + right
    return left + " " + right


def _should_join(prev, cur_rec, cur_geo, cur_page_idx):
    """Entscheidet, ob cur_rec die offene Zeile prev fortsetzt (C4-Kernregel).
    prev: {rec, page_idx, full, at_margin} - full/at_margin sind mit der
    Geometrie VON PREVS Seite berechnet (wichtig fuer den Seitenuebergang)."""
    ptext = prev["rec"]["text"].rstrip()
    # STARKES Textsignal, immer Pflicht: kein Satzschlusszeichen am Ende.
    if ptext and ptext[-1] in _JOIN_TERMINAL_PUNCT:
        return False
    # Geometrie: die obere Zeile muss voll sein (vom Setzer umbrochen).
    if not prev["full"]:
        return False
    # Aufzaehlungen bleiben zeilenweise (Marker oben oder unten -> kein Join).
    if _is_list_marker(prev["rec"]["text"]) or _is_list_marker(cur_rec["text"]):
        return False
    cs = cur_rec["text"].lstrip()
    if not cs:
        return False
    # Auf derselben Seite: gleiche Absatz-Block-Zugehoerigkeit (vertikaler
    # Abstand). Ueber Seitengrenzen gibt es kein Abstandssignal - dort traegt
    # die Entscheidung allein Geometrie (voll) + Text.
    if prev["page_idx"] == cur_page_idx:
        leading = cur_geo["leading"]
        if (leading and cur_rec.get("top") is not None
                and prev["rec"].get("top") is not None):
            if cur_rec["top"] - prev["rec"]["top"] > leading * _JOIN_BLOCK_FACTOR:
                return False
    # SCHWACHES Textsignal (nur Flattersatz): Folgezeile beginnt klein. Bei
    # echtem Blocksatz (prev auf geteiltem Justierungsrand) fallengelassen.
    if not prev["at_margin"]:
        first = cs[0]
        if first.isalpha() and first.isupper():
            return False
    return True


def _render_lines(pages_items, pages_tables=None):
    """Baut aus den markierten Zeilen-Records die Textzeilen und verbindet
    dabei umbrochene Fliesstext-Zeilen wieder zu Absaetzen (C4, Regel siehe
    oben). Headings und Pipe-Tabellen unveraendert wie in C2/C3.

    Leerzeilen: ein Absatz, in dem TATSAECHLICH verbunden wurde, bekommt
    Leerzeilen davor/dahinter; eine Zeile, die nichts gezogen hat, bleibt
    stehen wie zuvor. So werden Aufzaehlungen (Fixture d) und Adressbloecke
    (Fixture c) NICHT durch neue Leerzeilen aufgetrennt, waehrend verbundene
    Absaetze sauber getrennt sind. Ueberzaehlige Leerzeilen kollabiert
    _clean_text danach.

    Der offene Absatz (open_text) und prev laufen bewusst UEBER Seitengrenzen
    weiter - so kann ein Absatz, der unten auf Seite N voll umbricht, oben auf
    Seite N+1 fortgesetzt werden (Cross-Page-Join, Fixture h). Heading oder
    Tabelle schliessen ihn (flush + prev=None).

    Rueckgabe: (zeilen, stats) mit stats={"joins","tables"} fuer die
    Struktur-Note (wie viele Umbrueche verbunden, wie viele Tabellen emittiert)."""
    geos = [_page_geometry(page) for page in pages_items]
    doc_leading = _median([d for g in geos for d in g["deltas"]])
    for g in geos:
        if g["leading"] is None:
            g["leading"] = doc_leading

    out = []
    n_joins = 0           # verbundene Umbrueche (fuer die Struktur-Note)
    n_tables = 0          # emittierte Pipe-Tabellen (fuer die Struktur-Note)
    open_text = None      # aktuell offener Absatz (String) oder None
    open_joined = False   # wurde in diesem Absatz mindestens einmal verbunden?
    prev = None           # letzte Body-Zeile: {rec, page_idx, full, at_margin}

    def flush():
        nonlocal open_text, open_joined
        if open_text is None:
            return
        if open_joined:
            out.append("")
            out.append(open_text)
            out.append("")
        else:
            out.append(open_text)
        open_text, open_joined = None, False

    for p_idx, page in enumerate(pages_items):
        tables = pages_tables[p_idx] if pages_tables else []
        geo = geos[p_idx]
        emitted = set()
        for r in page:
            k = r.get("table")
            if k is not None:
                flush()
                prev = None
                if k not in emitted:
                    emitted.add(k)
                    n_tables += 1
                    out.append("")
                    out.extend(_pipe_table(tables[k]["rows"]))
                    out.append("")
                continue
            level = r.get("heading")
            if level:
                flush()
                prev = None
                out.append("")
                out.append("#" * level + " " + r["text"].strip())
                out.append("")
                continue
            # Body-Zeile: an den offenen Absatz anfuegen oder neu beginnen.
            if (prev is not None and open_text is not None
                    and _should_join(prev, r, geo, p_idx)):
                open_text = _merge_join(open_text, r["text"])
                open_joined = True
                n_joins += 1
            else:
                flush()
                open_text = r["text"]
            prev = {"rec": r, "page_idx": p_idx,
                    "full": _line_full(r, geo),
                    "at_margin": _line_at_margin(r, geo)}
    flush()
    return out, {"joins": n_joins, "tables": n_tables}


def _ocr_pages(path: str, page_numbers, lang="de"):
    """
    Wandelt AUSGEWAEHLTE Seiten eines PDFs per OCR (Texterkennung) in Text um.
    page_numbers: Liste von 1-basierten Seitennummern.
    lang: Sprache der Anleitung im Fehlerfall (die Meldung landet in der Note).

    Rueckgabe: (texte, fehler)
      - texte:  dict {seitennummer: erkannter_text} - kann bei einem Fehler
                auch teilweise gefuellt sein (bereits erkannte Seiten).
      - fehler: None bei Erfolg, sonst eine verstaendliche Fehlermeldung.

    Braucht das Programm 'Tesseract' auf dem Rechner sowie die Python-Pakete
    pdf2image und pytesseract. Fehlt etwas, geben wir eine klare Anleitung
    zurueck statt abzustuerzen.
    """
    # Pakete pruefen
    try:
        from pdf2image import convert_from_path
        import pytesseract
    except ImportError:
        return {}, _nt(lang, "ocr_needed")

    # Falls wir Tesseract am Standardpfad gefunden haben, pytesseract dorthin zeigen.
    # (Damit ist kein PATH-Eintrag noetig.)
    if _TESSERACT_PATH:
        pytesseract.pytesseract.tesseract_cmd = _TESSERACT_PATH

    # Tesseract-Programm pruefen (pytesseract ruft es intern auf)
    try:
        pytesseract.get_tesseract_version()
    except Exception:
        return {}, _nt(lang, "tesseract_missing")

    # Nur die angefragten Seiten in Bilder wandeln und erkennen.
    # 150 dpi ist ein guter Kompromiss aus Genauigkeit und Tempo.
    # poppler_path nur mitgeben, wenn wir ihn gefunden haben.
    convert_kwargs = {"dpi": 150}
    if _POPPLER_PATH:
        convert_kwargs["poppler_path"] = _POPPLER_PATH

    texts = {}
    for n in page_numbers:
        try:
            images = convert_from_path(path, first_page=n, last_page=n,
                                        **convert_kwargs)
        except Exception as e:
            # Meist fehlt 'poppler' (wird von pdf2image gebraucht)
            return texts, _nt(lang, "poppler_missing", n=n, e=e)
        if not images:
            texts[n] = ""
            continue
        # Sowohl Deutsch als auch Englisch versuchen (viele Dokumente sind gemischt)
        try:
            txt = pytesseract.image_to_string(images[0], lang="deu+eng")
        except Exception:
            # Falls das deutsche Sprachpaket fehlt, nur Englisch
            txt = pytesseract.image_to_string(images[0])
        texts[n] = txt

    return texts, None


class UnreadablePdfError(ValueError):
    """PDF, bei der KEINE Seite lesbar war (auch nicht per OCR).
    Die Exception-Meldung ist die verstaendliche Meldung fuers UI;
    'detail' traegt den rohen technischen Grund (z.B. die pdfminer-
    Exception der ersten kaputten Seite) fuer die Fehlersuche."""

    def __init__(self, message: str, detail: str = ""):
        super().__init__(message)
        self.detail = detail


# Ab wie vielen Zeichen extrahierten Texts gilt eine PDF-Seite als "Textseite"?
# Darunter ist es praktisch sicher eine Bild-/Scan-Seite (oder leer).
_PAGE_TEXT_MIN = 15

# Grobe, dokumentierte Groessenordnung: was eine PDF-Seite als BILD an
# Tokens kosten wuerde, wenn man sie direkt ans Modell gibt.
_IMG_TOKENS_PER_PAGE = 1500


def extract_pdf(path: str, lang: str = "de") -> dict:
    """PDF: Text -> Markdown. Wenn ueberwiegend Tabellen -> CSV.
    Seiten ohne Text (Scans/Bilder) werden EINZELN per OCR erkannt -
    auch in gemischten PDFs (z.B. Text-Deckblatt + gescannter Anhang)."""
    import pdfplumber

    page_texts = []       # extrahierter Text pro Seite (Index 0 = Seite 1)
    page_recs = []        # Zeilen-Records mit Geometrie pro Seite (oder None)
    page_tbl_infos = []   # pro Seite: [{"bbox","rows"}] fuer Pipe-Tabellen (C3)
    extra_texts = []      # Text AUSSERHALB von Tabellen pro Seite (fuer CSV-Zweig)
    all_tables = []
    n_images = 0
    broken_pages = []     # Seiten, deren Extraktion eine Exception warf
    first_error = ""      # technischer Grund der ersten kaputten Seite
    with pdfplumber.open(path) as pdf:
        for page_no, page in enumerate(pdf.pages, start=1):
            # JEDE Seite in ihrem eigenen try/except: eine einzelne kaputte
            # Seite (z.B. pdfminer-ValueError) darf nicht mehr die ganze
            # Datei mitreissen. Kaputte Seiten bleiben ohne Text und werden
            # dadurch unten wie Bildseiten behandelt (OCR-Rettung).
            try:
                txt = page.extract_text() or ""
                # Zeilen-Records fuer die Strukturlogik (Block C); None bei
                # Abweichung vom extract_text()-Ergebnis -> Seite laeuft
                # dann ohne Struktur-Extras weiter.
                recs = _line_records(page, txt)
                page_images = len(page.images or [])
                tables_on_page = page.find_tables()
                page_tables = []
                tbl_infos = []
                for t in tables_on_page:
                    data = t.extract()
                    if data:
                        page_tables.append(data)
                        # fuer den Markdown-Zweig (C3): Tabellen mit BBox,
                        # aber nur wenn mindestens EINE Zelle Inhalt hat
                        # (leere Gitter sind keine Tabellen, nur Linien)
                        if any(c is not None and str(c).strip()
                               for row in data for c in row):
                            tbl_infos.append({"bbox": t.bbox, "rows": data})
                # Begleittext ausserhalb der Tabellen-Bereiche merken, damit
                # der CSV-Zweig ihn nicht stillschweigend verwirft.
                # strict=False: reale PDFs enthalten Tabellen, deren Box durch
                # Rundungsartefakte um Bruchteile eines Punkts ueber den
                # Seitenrand ragt - pdfplumber (0.11.x) wirft dann per Default
                # ValueError statt still zu beschneiden.
                if tables_on_page:
                    region = page
                    for t in tables_on_page:
                        region = region.outside_bbox(t.bbox, strict=False)
                    extra = region.extract_text() or ""
                else:
                    extra = txt
            except Exception as e:
                broken_pages.append(page_no)
                if not first_error:
                    first_error = f"Seite {page_no}: {e}"
                txt, extra, page_tables, page_images = "", "", [], 0
                recs, tbl_infos = None, []
            page_texts.append(txt)
            page_recs.append(recs)
            page_tbl_infos.append(tbl_infos)
            extra_texts.append(extra)
            all_tables.extend(page_tables)
            n_images += page_images

    n_pages = len(page_texts)
    # raw_text = ungefilterte Extraktion (fuer Token-Vergleich "vorher")
    raw_text = "\n".join(page_texts)
    text_len = len(raw_text.strip())

    # --- Bild-Seiten-Erkennung: PRO SEITE statt global ---
    # Eine Seite mit praktisch keinem Text ist eine Bild-/Scan-Seite.
    image_pages = [i for i, t in enumerate(page_texts, start=1)
                   if len(t.strip()) < _PAGE_TEXT_MIN]
    text_pages = [i for i in range(1, n_pages + 1) if i not in image_pages]

    # Hinweis auf eingebettete Bilder in Text-PDFs (prinzipbedingt nicht
    # extrahierbar - aber wir sagen es, statt zu schweigen).
    image_note = ""
    if n_images > 0 and text_pages:
        image_note = _nt(lang, "image_note", n=n_images)

    # ---- Fall 1: reine Bild-PDF (alle Seiten ohne Text) -> komplette OCR ----
    if image_pages and not text_pages:
        ocr_texts, ocr_error = _ocr_pages(path, image_pages, lang=lang)
        if ocr_error and not any(t.strip() for t in ocr_texts.values()):
            if broken_pages and len(broken_pages) == n_pages:
                # ALLE Seiten warfen Exceptions und OCR konnte nichts retten
                # -> sauberer Fehler (HTTP 400) ueber den bestehenden Pfad,
                # mit verstaendlicher Meldung statt roher pdfminer-Exception.
                raise UnreadablePdfError(
                    _et(lang, "pdf_unreadable", n=n_pages),
                    detail=" | ".join(x for x in (first_error, ocr_error) if x))
            # OCR nicht moeglich -> ehrlich melden statt leeres Ergebnis
            note = _nt(lang, "bildpdf_no_text", err=ocr_error)
            if broken_pages:
                note = _nt(lang, "broken_skipped_prefix",
                           pages=", ".join(map(str, broken_pages))) + note
            return {
                "raw_text": "",
                "converted": "",
                "target_format": "markdown",
                "note": note,
            }
        ocr_text = "\n\n".join(ocr_texts.get(i, "") for i in image_pages)
        md = _clean_text(ocr_text)
        rescued_broken = [i for i in broken_pages
                          if ocr_texts.get(i, "").strip()]
        if broken_pages and len(broken_pages) == n_pages:
            # Keine echte Bild-PDF, sondern defekte Seiten - ehrlich benennen.
            note = _nt(lang, "rescued_all_pages", n=n_pages)
        else:
            note = _nt(lang, "bildpdf_ocr", n=n_pages)
            if rescued_broken:
                note += _nt(lang, "rescued_pages",
                            pages=", ".join(map(str, rescued_broken)))
        if ocr_error:
            missing = [i for i in image_pages if not ocr_texts.get(i, "").strip()]
            note = _nt(lang, "ocr_aborted",
                       pages=", ".join(map(str, missing)),
                       err=ocr_error, note=note)
        return {
            "raw_text": ocr_text,
            "converted": md,
            "target_format": "markdown",
            "note": note,
            "was_ocr": True,
            "n_pages": n_pages,
            # "Vorher" = was die Seiten als BILDER kosten wuerden (Schaetzung)
            "tokens_before_hint": n_pages * _IMG_TOKENS_PER_PAGE,
        }

    # ---- Fall 2: GEMISCHTES PDF (Textseiten + Bildseiten) ----
    if image_pages and text_pages:
        ocr_texts, ocr_error = _ocr_pages(path, image_pages, lang=lang)
        # Seiten in ORIGINAL-Reihenfolge zusammensetzen:
        # Textseiten direkt (mit Zeilen-Records fuer die Strukturlogik),
        # Bildseiten aus der OCR (nur Text, keine Geometrie).
        assembled = []
        pages_items = []
        pages_tbl = []   # Pipe-Tabellen nur fuer Seiten mit Geometrie-Records
        for i in range(1, n_pages + 1):
            if i in ocr_texts and ocr_texts[i].strip():
                assembled.append(ocr_texts[i])
                pages_items.append(_plain_records(ocr_texts[i]))
                pages_tbl.append([])
            elif i in image_pages:
                assembled.append("")  # OCR fehlgeschlagen -> Seite fehlt
                pages_items.append(_plain_records(""))
                pages_tbl.append([])
            else:
                assembled.append(page_texts[i - 1])
                recs = page_recs[i - 1]
                if recs is not None:
                    pages_items.append(recs)
                    pages_tbl.append(page_tbl_infos[i - 1])
                else:
                    pages_items.append(_plain_records(page_texts[i - 1]))
                    pages_tbl.append([])

        # WICHTIG: Kopf-/Fusszeilen-Entfernung laeuft VOR der Strukturlogik
        # (zeilenbasierte Wiederholungserkennung braucht die Originalzeilen).
        cleaned_pages, removed, removed_patterns = \
            _strip_repeating_headers_footers(pages_items,
                                             key=lambda r: r["text"])
        _mark_table_zones(cleaned_pages, pages_tbl)
        n_headings = _mark_headings(cleaned_pages)
        lines, sstats = _render_lines(cleaned_pages, pages_tbl)
        md = _clean_text("\n".join(lines))
        struct = _structure_note(n_headings, sstats["tables"], sstats["joins"],
                                 lang=lang)

        ocr_ok = [i for i in image_pages if ocr_texts.get(i, "").strip()]
        ocr_failed = [i for i in image_pages if i not in ocr_ok]
        # Defekte Seiten (Exception im Extraktions-Loop) von echten
        # Bild-Seiten unterscheiden - die Note soll nicht luegen.
        rescued_broken = [i for i in ocr_ok if i in broken_pages]
        image_ok = [i for i in ocr_ok if i not in broken_pages]
        if ocr_failed:
            failed_image = [i for i in ocr_failed if i not in broken_pages]
            failed_broken = [i for i in ocr_failed if i in broken_pages]
            parts = []
            if failed_image:
                parts.append(_nt(lang, "mixed_image_fail",
                                 pages=", ".join(map(str, failed_image))))
            if failed_broken:
                parts.append(_nt(lang, "mixed_broken_fail",
                                 pages=", ".join(map(str, failed_broken))))
            note = _nt(lang, "mixed_incomplete", parts=" ".join(parts),
                       err=(ocr_error or ""), n=len(text_pages))
        else:
            note = _nt(lang, "mixed_ok", n=len(text_pages))
            if image_ok:
                note += _nt(lang, "mixed_ocr_pages",
                            pages=", ".join(map(str, image_ok)))
        if struct:
            note += " " + struct
        if rescued_broken:
            note += _nt(lang, "rescued_pages",
                        pages=", ".join(map(str, rescued_broken)))
        if removed > 0:
            note += (_nt(lang, "headers_removed", n=removed)
                     + _removed_note(removed_patterns, lang=lang))

        # Korrigierte Token-Rechnung: Textseiten zaehlen als Text,
        # Bildseiten als geschaetzte Bild-Kosten.
        text_raw = "\n".join(page_texts[i - 1] for i in text_pages)
        before_hint = (count_tokens(text_raw)["count"]
                       + len(image_pages) * _IMG_TOKENS_PER_PAGE)
        return {
            "raw_text": "\n".join(assembled),
            "converted": md,
            "target_format": "markdown",
            "note": note,
            "was_ocr": True,
            "n_pages": n_pages,
            "tokens_before_hint": before_hint,
        }

    # ---- Fall 3: normale Text-PDF ----
    pages_items = []
    pages_tbl = []   # Pipe-Tabellen nur fuer Seiten mit Geometrie-Records
    for i, t in enumerate(page_texts):
        if page_recs[i] is not None:
            pages_items.append(page_recs[i])
            pages_tbl.append(page_tbl_infos[i])
        else:
            pages_items.append(_plain_records(t))
            pages_tbl.append([])
    # bereinigt: wiederkehrende Kopf-/Fusszeilen raus. Laeuft VOR der
    # Strukturlogik - die Wiederholungserkennung braucht die Originalzeilen.
    cleaned_pages, removed, removed_patterns = \
        _strip_repeating_headers_footers(pages_items, key=lambda r: r["text"])
    # Grobe Heuristik: viele Tabellenzeilen + wenig Fliesstext -> Tabellen-PDF
    table_cells = sum(len(r) for t in all_tables for r in t)

    if all_tables and table_cells > 0 and text_len < table_cells * 6:
        # Als Tabellen-Dokument behandeln -> CSV
        # Alle Tabellen untereinander, durch Leerzeile getrennt
        csv_blocks = [_rows_to_csv(t) for t in all_tables]
        converted = "\n\n".join(csv_blocks)
        note = _nt(lang, "table_pdf")
        # DEFENSIV: Fliesstext ausserhalb der Tabellen darf nicht
        # stillschweigend verloren gehen - er kommt als Begleittext dazu.
        companion = _clean_text("\n".join(extra_texts))
        if companion:
            converted = (f"# Begleittext (ausserhalb der Tabellen):\n"
                         f"{companion}\n\n{converted}")
            note += _nt(lang, "companion_kept")
        return {
            "raw_text": raw_text,
            "converted": converted,
            "target_format": "csv",
            "note": note + image_note,
        }

    # Standardfall: Fliesstext -> Markdown (kopf-/fusszeilenbereinigt,
    # konservativ erkannte Ueberschriften als #/##/###, Gitter-Tabellen
    # als Pipe-Tabellen inline im Textfluss)
    _mark_table_zones(cleaned_pages, pages_tbl)
    n_headings = _mark_headings(cleaned_pages)
    lines, sstats = _render_lines(cleaned_pages, pages_tbl)
    md = _clean_text("\n".join(lines))
    struct = _structure_note(n_headings, sstats["tables"], sstats["joins"],
                             lang=lang)
    if removed > 0 or struct:
        note = _nt(lang, "pdf_cleaned")
        if struct:
            note += " " + struct
        if removed > 0:
            note += (_nt(lang, "headers_removed", n=removed)
                     + _removed_note(removed_patterns, lang=lang))
    else:
        note = _nt(lang, "pdf_plain")
    return {
        "raw_text": raw_text,
        "converted": md,
        "target_format": "markdown",
        "note": note + image_note,
    }


def extract_docx(path: str, lang: str = "de") -> dict:
    """Word: Ueberschriften -> Markdown-Headings, Tabellen -> Markdown-Tabellen.
    Wichtig: Absaetze und Tabellen bleiben in ORIGINAL-Reihenfolge
    (frueher landeten alle Tabellen gesammelt am Ende - Bezuege wie
    'die folgende Tabelle' waren damit kaputt)."""
    import docx
    from docx.oxml.text.paragraph import CT_P
    from docx.oxml.table import CT_Tbl
    from docx.text.paragraph import Paragraph
    from docx.table import Table

    doc = docx.Document(path)
    md_parts = []
    raw_parts = []

    def render_paragraph(para):
        text = para.text.strip()
        raw_parts.append(para.text)
        if not text:
            return
        # para.style ist None, wenn styles.xml keinen Default-Paragraph-Style
        # definiert (Nicht-Word-Erzeuger wie docx-js oder pandoc-Minimal).
        # Dann als Fliesstext behandeln, keine Ueberschriften-Erkennung.
        style = (para.style.name or "").lower() if para.style is not None else ""
        if style.startswith("heading 1") or style == "title":
            md_parts.append(f"# {text}")
        elif style.startswith("heading 2"):
            md_parts.append(f"## {text}")
        elif style.startswith("heading 3"):
            md_parts.append(f"### {text}")
        elif style.startswith("list"):
            md_parts.append(f"- {text}")
        else:
            md_parts.append(text)

    def render_table(table):
        # Wichtig: die Tabellenzeilen muessen DIREKT untereinander stehen
        # (ein Block mit einfachen Zeilenumbruechen), sonst ist die Markdown-
        # Tabelle ungueltig. Deshalb bauen wir jede Tabelle als einen String.
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        raw_parts.append("\n".join("\t".join(r) for r in rows))
        if not rows:
            return
        header = rows[0]
        table_lines = []
        table_lines.append("| " + " | ".join(header) + " |")
        table_lines.append("| " + " | ".join(["---"] * len(header)) + " |")
        for r in rows[1:]:
            # auf Header-Breite normalisieren
            r = (r + [""] * len(header))[: len(header)]
            table_lines.append("| " + " | ".join(r) + " |")
        md_parts.append("\n".join(table_lines))

    # Ueber die Body-Kinder in Dokument-Reihenfolge laufen:
    # CT_P = Absatz, CT_Tbl = Tabelle. Andere Elemente (z.B. Abschnitts-
    # eigenschaften) sind kein Inhalt und werden uebersprungen.
    for child in doc.element.body.iterchildren():
        if isinstance(child, CT_P):
            render_paragraph(Paragraph(child, doc))
        elif isinstance(child, CT_Tbl):
            render_table(Table(child, doc))

    converted = _clean_text("\n\n".join(md_parts))
    raw_text = "\n".join(raw_parts)
    return {
        "raw_text": raw_text,
        "converted": converted,
        "target_format": "markdown",
        "note": _nt(lang, "docx_note"),
    }


def extract_xlsx(path: str, lang: str = "de") -> dict:
    """Excel: jedes Blatt -> CSV. Daten gehoeren in CSV (am wenigsten Tokens).

    Formel-Zellen: Wir lesen das Workbook DOPPELT - einmal mit data_only=True
    (liefert von Excel zwischengespeicherte Rechenergebnisse) und einmal mit
    data_only=False (liefert die Formel-Strings). Bei programmatisch erzeugten
    Dateien existiert kein Ergebnis-Cache; frueher wurden solche Zellen zu
    LEEREN Feldern. Jetzt geben wir stattdessen die Formel selbst aus."""
    import openpyxl

    wb_val = openpyxl.load_workbook(path, data_only=True, read_only=True)
    wb_formula = openpyxl.load_workbook(path, data_only=False, read_only=True)
    blocks = []
    raw_parts = []
    formula_fallbacks = 0

    for ws_val, ws_formula in zip(wb_val.worksheets, wb_formula.worksheets):
        rows = []
        formula_rows = [list(r) for r in ws_formula.iter_rows(values_only=True)]
        for r_idx, row in enumerate(ws_val.iter_rows(values_only=True)):
            cells = list(row) if row is not None else []
            f_row = formula_rows[r_idx] if r_idx < len(formula_rows) else []
            merged = []
            for c_idx, val in enumerate(cells):
                if val is None and c_idx < len(f_row):
                    f_val = f_row[c_idx]
                    # Zelle hat kein zwischengespeichertes Ergebnis, aber
                    # eine Formel -> Formel-String ausgeben statt Leere.
                    if isinstance(f_val, str) and f_val.startswith("="):
                        merged.append(f_val)
                        formula_fallbacks += 1
                        continue
                merged.append(val)
            # komplett leere Zeilen ueberspringen
            if not merged or all(c is None for c in merged):
                continue
            rows.append(merged)
        if not rows:
            continue
        csv_text = _rows_to_csv(rows)
        raw_parts.append(csv_text)
        # Blattname als Kommentarzeile, damit man mehrere Blaetter unterscheidet
        if len(wb_val.worksheets) > 1:
            blocks.append(f"# Blatt: {ws_val.title}\n{csv_text}")
        else:
            blocks.append(csv_text)

    converted = "\n\n".join(blocks).strip()
    raw_text = "\n\n".join(raw_parts)
    note = _nt(lang, "xlsx_note")
    if formula_fallbacks > 0:
        note += _nt(lang, "xlsx_formula", n=formula_fallbacks)
    return {
        "raw_text": raw_text,
        "converted": converted,
        "target_format": "csv",
        "note": note,
    }


def extract_csv(path: str, lang: str = "de") -> dict:
    """CSV ist schon optimal - nur einlesen und leicht saeubern."""
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        content = f.read()
    cleaned = content.replace("\r\n", "\n").replace("\r", "\n").strip()
    return {
        "raw_text": content,
        "converted": cleaned,
        "target_format": "csv",
        "note": _nt(lang, "csv_note"),
    }


def extract_txt(path: str, lang: str = "de") -> dict:
    """Reiner Text -> als Markdown bereinigen."""
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        content = f.read()
    cleaned = _clean_text(content)
    return {
        "raw_text": content,
        "converted": cleaned,
        "target_format": "markdown",
        "note": _nt(lang, "txt_note"),
    }


def extract_pptx(path: str, lang: str = "de") -> dict:
    """PowerPoint via python-pptx -> strukturiertes Markdown.

    Pro Folie: Titel als Trenner "## Folie N: <Titel>", Body-Text als
    verschachtelte Liste (Ebene ueber Einrueckung), Untertitel plain (keine
    erfundene Aufzaehlung), Tabellen als Markdown-Pipe-Tabellen,
    Sprechernotizen unter "### Notes:".

    Frueher lief das ueber markitdown. Das zog magika als schwergewichtige
    Abhaengigkeit nach, deren ONNX-Model-Verzeichnis im eingefrorenen
    Zustand (.exe) fehlte und PPTX zum Absturz brachte. python-pptx ist
    ohnehin Dependency und findet den Dateityp bereits ueber die Endung.
    """
    from pptx import Presentation
    from pptx.enum.shapes import PP_PLACEHOLDER

    prs = Presentation(path)

    def _table_md(table):
        rows = ["| " + " | ".join(c.text.strip() for c in r.cells) + " |"
                for r in table.rows]
        if rows:
            rows.insert(1, "| " + " | ".join(["---"] * len(table.columns)) + " |")
        return "\n".join(rows)

    def _is_subtitle(shape):
        return (shape.is_placeholder
                and shape.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE)

    def _body_md(shape):
        # Untertitel plain, sonst Liste mit Einrueckung nach Absatz-Ebene
        plain = _is_subtitle(shape)
        lines = []
        for para in shape.text_frame.paragraphs:
            t = para.text.strip()
            if not t:
                continue
            lines.append(t if plain else "  " * para.level + "- " + t)
        return "\n".join(lines)

    blocks = []
    n_notes = 0
    for i, slide in enumerate(prs.slides, 1):
        title_shape = slide.shapes.title
        # Titel ueber shape_id erkennen - slide.shapes.title liefert ein
        # anderes Wrapper-Objekt als die Iteration, "is" greift nicht.
        title_id = title_shape.shape_id if title_shape is not None else None
        title = title_shape.text.strip() if title_shape is not None else ""

        part = [f"## Folie {i}: {title}" if title else f"## Folie {i}"]
        for shape in slide.shapes:
            if title_id is not None and shape.shape_id == title_id:
                continue  # Titel steckt schon im Trenner
            if shape.has_table:
                part.append(_table_md(shape.table))
            elif shape.has_text_frame and shape.text_frame.text.strip():
                body = _body_md(shape)
                if body:
                    part.append(body)

        if slide.has_notes_slide:
            note_text = (slide.notes_slide.notes_text_frame.text or "").strip()
            if note_text:
                n_notes += 1
                part.append("### Notes:\n" + note_text)

        blocks.append("\n".join(part))

    text = "\n\n".join(blocks)
    cleaned = _clean_text(text)

    note = _nt(lang, "pptx_note")
    # Fairness-Hinweis: Sprechernotizen landen mit im Output - das ist je
    # nach Datei gewollt oder ueberraschend, deshalb sagen wir es dazu.
    if n_notes > 0:
        note += _nt(lang, "pptx_speaker", n=n_notes)

    return {
        "raw_text": text,
        "converted": cleaned,
        "target_format": "markdown",
        "note": note,
    }


# ---------------------------------------------------------------------------
# MODELLSPEZIFISCHES WRAPPING
# ---------------------------------------------------------------------------

def wrap_for_model(content: str, source_name: str, fmt: str, model: str) -> str:
    """
    Verpackt den Inhalt so, wie es das jeweilige Ziel-Modell am liebsten mag:
      claude -> XML-Tags (offizielle Anthropic-Empfehlung)
      gpt    -> Markdown-Ueberschrift + Triple-Quote-Delimiter (OpenAI-Stil)
      gemini -> schlichte Markdown-Struktur mit Label
      none   -> unveraendert
    Alle Varianten kosten nur wenige Tokens und helfen dem Modell, das
    Dokument sauber von der eigentlichen Aufgabe zu trennen.
    """
    safe_name = os.path.basename(source_name)

    if model == "claude":
        return (f'<document source="{safe_name}" format="{fmt}">\n'
                f"{content}\n"
                f"</document>")

    if model == "gpt":
        return (f'## Dokument: {safe_name} (Format: {fmt})\n'
                f'"""\n{content}\n"""')

    if model == "gemini":
        return (f'**Dokument: {safe_name}** (Format: {fmt})\n\n'
                f'{content}\n\n'
                f'--- Ende des Dokuments ---')

    return content  # "none" oder unbekannt -> pur


def wrap_xml(content: str, source_name: str, fmt: str) -> str:
    """Rueckwaerts-kompatibler Alias (alte Aufrufe -> Claude-Stil)."""
    return wrap_for_model(content, source_name, fmt, "claude")


# ---------------------------------------------------------------------------
# HAUPT-DISPATCHER
# ---------------------------------------------------------------------------

_EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
    ".xlsm": extract_xlsx,
    ".csv": extract_csv,
    ".txt": extract_txt,
    ".md": extract_txt,
    ".pptx": extract_pptx,
}

SUPPORTED_EXTENSIONS = sorted(_EXTRACTORS.keys())


def convert_file(path: str, add_xml: bool = False, target_model: str = None,
                 original_name: str = None, lang: str = "de") -> dict:
    """
    Haupteinstieg. Nimmt einen Dateipfad, gibt ein Ergebnis-dict zurueck.
    target_model:  'claude' | 'gpt' | 'gemini' | 'none'
                   (steuert das modellspezifische Wrapping)
    add_xml:       alter Parameter, bleibt aus Kompatibilitaet -
                   True wirkt wie target_model='claude'.
    original_name: echter Dateiname fuers Wrapping (falls der Pfad nur ein
                   temporaerer Name ist, z.B. beim Server-Upload).
    lang:          Sprache der Anzeige-Note ('de'/'en', sonst 'de').
      tokens_saved:  Differenz
      percent_saved: Ersparnis in Prozent
      token_method:  'tiktoken' oder 'estimate'
    """
    lang = _norm_lang(lang)
    ext = os.path.splitext(path)[1].lower()
    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        return {
            "ok": False,
            "error": _et(lang, "unsupported_type", ext=ext,
                         exts=", ".join(SUPPORTED_EXTENSIONS)),
        }

    try:
        result = extractor(path, lang=lang)
    except UnreadablePdfError as e:
        # Verstaendliche Hauptmeldung; der rohe technische Grund kommt
        # separat mit (fuers UI als ausklappbares Detail, nicht als
        # Hauptfehler).
        return {"ok": False, "error": str(e), "error_detail": e.detail}
    except Exception as e:
        return {"ok": False, "error": _et(lang, "read_failed", e=e)}

    source_name = os.path.basename(original_name) if original_name else os.path.basename(path)
    converted = result["converted"]
    fmt = result["target_format"]

    # --- Ehrliche Token-Messung ---
    # "vorher": was du zahlst, wenn du die rohe Extraktion ungefiltert ans
    #           LLM gibst (inkl. Extraktions-Muell, doppelte Umbrueche etc.)
    # "nachher": der bereinigte, optimierte Inhalt OHNE XML-Deko.
    after_clean = count_tokens(converted)

    # Sonderfall OCR (Bild-PDF oder gemischtes PDF): Hier liefert der
    # Extraktor einen ehrlichen "Vorher"-Wert mit (tokens_before_hint):
    #   - reine Bild-PDF: Seiten x ~1500 Tokens (Kosten als Bild)
    #   - gemischtes PDF: Text der Textseiten + Bildseiten x ~1500 Tokens
    # Bild-Kosten sind immer eine Schaetzung, das weisen wir transparent aus.
    is_ocr = result.get("was_ocr", False)
    if "tokens_before_hint" in result:
        before_count = result["tokens_before_hint"]
        before_method = "estimate"
    elif is_ocr:
        # Fallback fuer Extraktoren ohne Hint (sollte nicht vorkommen)
        n_pages = result.get("n_pages", 1)
        before_count = n_pages * _IMG_TOKENS_PER_PAGE
        before_method = "estimate"
    else:
        before = count_tokens(result["raw_text"])
        before_count = before["count"]
        before_method = before["method"]

    saved = before_count - after_clean["count"]
    # OCR-Faelle: Bild-Seiten sind fuer Textmodelle nicht lesbar, ein
    # Prozentvergleich waere Pseudo-Genauigkeit -> kein percent_saved;
    # token_method beschreibt dann die Zaehlung des Ergebnisses.
    # tokens_before bleibt als Schaetzgrundlage (Statistik) erhalten.
    if is_ocr:
        percent = None
        method_out = after_clean["method"]
    else:
        percent = round((saved / before_count * 100) if before_count > 0 else 0.0, 1)
        method_out = before_method

    # Ziel-Modell bestimmen (add_xml=True bleibt als Alias fuer Claude)
    if target_model is None:
        target_model = "claude" if add_xml else "none"
    if target_model not in ("claude", "gpt", "gemini", "none"):
        target_model = "none"

    # Finale Ausgabe (ggf. modellspezifisch verpackt)
    output_text = converted
    wrap_overhead = 0
    if target_model != "none":
        output_text = wrap_for_model(converted, source_name, fmt, target_model)
        wrap_overhead = count_tokens(output_text)["count"] - after_clean["count"]

    out_ext = ".csv" if fmt == "csv" else ".md"

    return {
        "ok": True,
        "source_name": source_name,
        "target_format": fmt,
        "note": result["note"],
        "output_text": output_text,
        "output_ext": out_ext,
        "tokens_before": before_count,
        "tokens_after": after_clean["count"],   # bereinigt, ohne Wrapping
        "tokens_saved": saved,
        "percent_saved": percent,
        "token_method": method_out,
        "was_ocr": is_ocr,
        "target_model": target_model,
        "wrap_overhead": wrap_overhead,          # was das Wrapping extra kostet
        # alte Feldnamen fuer Kompatibilitaet:
        "xml_wrapped": target_model != "none",
        "xml_overhead": wrap_overhead,
        "tokens_final": count_tokens(output_text)["count"],  # inkl. Wrapping
    }
