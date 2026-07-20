# -*- coding: utf-8 -*-
"""
test_block1.py
--------------
Regressionstests für Block 1 (Konverter-Fixes). Jeder Test konstruiert seine
Testdatei selbst (reportlab / python-docx / openpyxl / python-pptx) und
beweist damit den jeweiligen Fix:

  A) Gemischte PDFs: Textseiten + Bildseiten -> Pro-Seite-OCR, nichts fehlt
  B) Kapitelüberschriften überleben das Header-Stripping (EDGE + _PROTECTED)
  D) DOCX: Absätze und Tabellen in Original-Reihenfolge
  E) XLSX: Formeln ohne Ergebnis-Cache -> Formel-String statt leerer Zelle
  Nebenpunkte: PPTX-Notes-Hinweis, defensiver CSV-Zweig (Begleittext),
  Bilder-Hinweis, korrigierte Token-Rechnung bei gemischten PDFs.

Aufruf:  python tests/test_block1.py
"""

import os
import sys
import traceback

# converter.py liegt eine Ebene höher
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from converter import convert_file, count_tokens  # noqa: E402

FIXTURES = os.path.join(os.path.dirname(os.path.abspath(__file__)), "fixtures")
os.makedirs(FIXTURES, exist_ok=True)

A4_PT = (595.27, 841.89)


# ---------------------------------------------------------------------------
# FIXTURE-BAUER
# ---------------------------------------------------------------------------

def _make_text_image(lines, size=(1240, 1754), font_size=54):
    # font_size 54: gross genug fuer sichere OCR, klein genug, dass auch
    # lange Zeilen nicht ueber den rechten Bildrand hinauslaufen.
    """Erzeugt ein weißes Bild mit großem schwarzen Text (gut für OCR)."""
    from PIL import Image, ImageDraw, ImageFont
    img = Image.new("RGB", size, "white")
    draw = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("arial.ttf", font_size)
    except OSError:
        font = ImageFont.load_default()
    y = 200
    for line in lines:
        draw.text((100, y), line, fill="black", font=font)
        y += font_size + 60
    return img


def build_mixed_pdf(path):
    """Seite 1: echter Text (Deckblatt). Seiten 2+3: nur Bilder mit Text."""
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader

    c = canvas.Canvas(path, pagesize=A4_PT)
    # Seite 1: Text-Deckblatt
    c.setFont("Helvetica-Bold", 20)
    c.drawString(72, 760, "Quartalsbericht der Beispiel GmbH")
    c.setFont("Helvetica", 12)
    c.drawString(72, 720, "Dieses Deckblatt enthaelt echten, extrahierbaren Text.")
    c.drawString(72, 700, "Die folgenden Seiten sind gescannte Bildseiten.")
    c.showPage()
    # Seite 2 + 3: reine Bildseiten
    img2 = _make_text_image(["GEHEIME KENNZAHL EINS: 4711"])
    c.drawImage(ImageReader(img2), 0, 0, width=A4_PT[0], height=A4_PT[1])
    c.showPage()
    img3 = _make_text_image(["GEHEIME KENNZAHL ZWEI: 1337"])
    c.drawImage(ImageReader(img3), 0, 0, width=A4_PT[0], height=A4_PT[1])
    c.showPage()
    c.save()


def build_chapter_pdf(path):
    """5 Seiten mit Kopfzeile, Kapitelüberschrift oben, Fußzeile 'Seite N',
    und einer wiederholten kurzen Zeile MITTEN im Inhalt (darf nicht weg)."""
    from reportlab.pdfgen import canvas

    c = canvas.Canvas(path, pagesize=A4_PT)
    for n in range(1, 6):
        c.setFont("Helvetica", 11)
        c.drawString(72, 800, "Vertraulich - Beispiel GmbH")          # Kopfzeile
        c.setFont("Helvetica-Bold", 16)
        c.drawString(72, 770, f"Kapitel {n}")                          # geschützt!
        c.setFont("Helvetica", 12)
        c.drawString(72, 730, f"Erster Inhaltssatz auf Seite {n} mit eigenem Wortlaut.")
        c.drawString(72, 710, f"Zweiter Inhaltssatz: Analyse Nummer {n * 11}.")
        c.drawString(72, 690, f"Dritter Inhaltssatz: Ergebnisse des Bereichs {n * 7}.")
        c.drawString(72, 670, "Interner Vermerk 99")                   # Mitte: bleibt!
        c.drawString(72, 650, f"Vierter Inhaltssatz: Bewertung der Lage {n + 100}.")
        c.drawString(72, 630, f"Fuenfter Inhaltssatz: Ausblick fuer Quartal {n}.")
        c.setFont("Helvetica", 10)
        c.drawString(72, 60, f"Seite {n} von 5")                       # Fußzeile
        c.showPage()
    c.save()


def build_order_docx(path):
    """Absatz -> Tabelle -> Absatz. Die Reihenfolge ist der Testgegenstand."""
    import docx
    doc = docx.Document()
    doc.add_paragraph("Unsere Preisliste folgt in der naechsten Tabelle:")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Produkt"
    table.cell(0, 1).text = "Preis"
    table.cell(1, 0).text = "Widget"
    table.cell(1, 1).text = "9,99 EUR"
    doc.add_paragraph("Diese Preise gelten ab sofort.")
    doc.save(path)


def build_formula_xlsx(path):
    """Formelzelle ohne Ergebnis-Cache (Datei nie in Excel geöffnet)."""
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = 10
    ws["A2"] = 20
    ws["A3"] = "=SUM(A1:A2)"
    ws["B1"] = "Umsatz"
    wb.save(path)


def build_notes_pptx(path):
    """Folie mit Sprechernotiz."""
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Praesentationstitel Alpha"
    slide.placeholders[1].text = "Untertitel Beta"
    slide.notes_slide.notes_text_frame.text = "Geheime Sprechernotiz 555"
    prs.save(path)


def build_table_pdf_with_text(path):
    """PDF, das überwiegend aus einer Tabelle besteht, plus ein kurzer
    Begleitsatz, der früher im CSV-Zweig stillschweigend verloren ging."""
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib import colors

    doc = SimpleDocTemplate(path, pagesize=A4)
    styles = getSampleStyleSheet()
    data = [[f"R{r}C{c}" for c in range(1, 6)] for r in range(1, 11)]
    table = Table(data)
    table.setStyle(TableStyle([
        ("GRID", (0, 0), (-1, -1), 0.5, colors.black),
        ("FONTSIZE", (0, 0), (-1, -1), 8),
    ]))
    doc.build([
        table,
        Paragraph("Alle Preise verstehen sich netto.", styles["Normal"]),
    ])


def build_image_in_text_pdf(path):
    """Text-PDF mit eingebettetem Bild -> Bilder-Hinweis muss erscheinen."""
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader
    from PIL import Image

    c = canvas.Canvas(path, pagesize=A4_PT)
    c.setFont("Helvetica", 12)
    c.drawString(72, 780, "Dieser Bericht enthaelt ein wichtiges Diagramm unten.")
    c.drawString(72, 760, "Der Text drumherum ist normal extrahierbar.")
    diagram = Image.new("RGB", (300, 200), "lightblue")
    c.drawImage(ImageReader(diagram), 72, 480, width=200, height=130)
    c.showPage()
    c.save()


def build_plain_text_pdf(path):
    """Normale 3-seitige Text-PDF mit Kopf-/Fußzeilen (Regression Standardfall)."""
    from reportlab.pdfgen import canvas
    c = canvas.Canvas(path, pagesize=A4_PT)
    for n in range(1, 4):
        c.setFont("Helvetica", 10)
        c.drawString(72, 800, "Firmenname AG - Interner Bericht")
        c.setFont("Helvetica", 12)
        c.drawString(72, 760, f"Inhaltlicher Absatz Nummer {n} mit individuellem Text.")
        c.drawString(72, 740, f"Weitere Details zur Auswertung {n * 3} folgen hier.")
        c.setFont("Helvetica", 10)
        c.drawString(72, 60, f"Seite {n}")
        c.showPage()
    c.save()


# ---------------------------------------------------------------------------
# TESTS
# ---------------------------------------------------------------------------

def test_a_mixed_pdf_per_page_ocr():
    """Fix A: Gemischtes PDF verliert keine Bildseiten mehr."""
    path = os.path.join(FIXTURES, "gemischt.pdf")
    build_mixed_pdf(path)
    r = convert_file(path)
    assert r["ok"], f"Konvertierung fehlgeschlagen: {r.get('error')}"
    out = r["output_text"]
    # Textseite muss da sein
    assert "Quartalsbericht" in out, "Text-Deckblatt fehlt im Output"
    # OCR-Inhalte der Bildseiten muessen da sein (frueher: komplett weg)
    assert "4711" in out, "OCR-Inhalt von Bildseite 2 fehlt (Kennzahl 4711)"
    assert "1337" in out, "OCR-Inhalt von Bildseite 3 fehlt (Kennzahl 1337)"
    # Original-Reihenfolge: Deckblatt vor Seite 2 vor Seite 3
    assert out.index("Quartalsbericht") < out.index("4711") < out.index("1337"), \
        "Seiten-Reihenfolge stimmt nicht"
    # Transparente Note
    assert "OCR" in r["note"], f"Note nennt OCR nicht: {r['note']}"
    assert r["was_ocr"] is True
    # Korrigierte Token-Rechnung: 2 Bildseiten x 1500 + Text des Deckblatts
    text_tokens = count_tokens("Quartalsbericht der Beispiel GmbH")["count"]
    assert r["tokens_before"] >= 2 * 1500 + text_tokens, \
        f"tokens_before ({r['tokens_before']}) enthaelt die Bildseiten-Schaetzung nicht"
    # Block G1: OCR-Pfade liefern keinen Pseudo-Prozentvergleich mehr;
    # token_method beschreibt die Zaehlung des ERGEBNISSES.
    assert r["percent_saved"] is None, \
        f"OCR darf keinen Pseudo-Prozentwert liefern: {r['percent_saved']}"
    assert r["token_method"] == count_tokens("probe")["method"], \
        "token_method muss die Ergebnis-Zaehlung beschreiben"


def test_b_chapter_headings_protected():
    """Fix B: Kapitelüberschriften überleben, Kopf-/Fußzeilen fliegen raus."""
    path = os.path.join(FIXTURES, "kapitel.pdf")
    build_chapter_pdf(path)
    r = convert_file(path)
    assert r["ok"], f"Konvertierung fehlgeschlagen: {r.get('error')}"
    out = r["output_text"]
    # Alle 5 Kapitelueberschriften muessen ueberleben (frueher: alle geloescht)
    for n in range(1, 6):
        assert f"Kapitel {n}" in out, f"'Kapitel {n}' wurde faelschlich entfernt!"
    # Echte Kopf-/Fusszeilen muessen weiterhin entfernt werden
    assert "Vertraulich - Beispiel GmbH" not in out, "Kopfzeile wurde nicht entfernt"
    assert "Seite 1 von 5" not in out, "Fusszeile wurde nicht entfernt"
    # Wiederholte Zeile MITTEN im Inhalt darf nicht entfernt werden (EDGE=2)
    assert "Interner Vermerk 99" in out, \
        "Zeile in Seitenmitte wurde entfernt - EDGE-Schutz greift nicht"
    # Transparenz: entfernte Zeilen stehen in der Note
    assert "Entfernt:" in r["note"], f"Note zeigt entfernte Zeilen nicht: {r['note']}"
    assert "Vertraulich" in r["note"], "Entfernte Kopfzeile fehlt in der Note"


def test_d_docx_original_order():
    """Fix D: Absatz -> Tabelle -> Absatz bleibt in dieser Reihenfolge."""
    path = os.path.join(FIXTURES, "reihenfolge.docx")
    build_order_docx(path)
    r = convert_file(path)
    assert r["ok"], f"Konvertierung fehlgeschlagen: {r.get('error')}"
    out = r["output_text"]
    i_before = out.index("Preisliste folgt")
    i_table = out.index("Widget")
    i_after = out.index("gelten ab sofort")
    assert i_before < i_table, "Tabelle steht VOR ihrem Einleitungssatz"
    assert i_table < i_after, "Tabelle wurde ans Ende verschoben (alter Bug)"
    # Markdown-Tabelle korrekt gebaut
    assert "| Produkt | Preis |" in out, "Markdown-Tabellenkopf fehlt"


def test_e_xlsx_formula_fallback():
    """Fix E: Formel ohne Cache -> Formel-String statt leerer Zelle."""
    path = os.path.join(FIXTURES, "formeln.xlsx")
    build_formula_xlsx(path)
    r = convert_file(path)
    assert r["ok"], f"Konvertierung fehlgeschlagen: {r.get('error')}"
    out = r["output_text"]
    assert "=SUM(A1:A2)" in out, "Formel-Zelle ist leer statt Formel-String (alter Bug)"
    assert "10" in out and "20" in out and "Umsatz" in out, "Normale Werte fehlen"
    assert "Formel" in r["note"], f"Note erklaert den Formel-Fallback nicht: {r['note']}"


def test_pptx_notes_hint():
    """Nebenpunkt: Sprechernotizen werden extrahiert UND in der Note erwähnt."""
    path = os.path.join(FIXTURES, "notizen.pptx")
    build_notes_pptx(path)
    r = convert_file(path)
    assert r["ok"], f"Konvertierung fehlgeschlagen: {r.get('error')}"
    out = r["output_text"]
    assert "Praesentationstitel Alpha" in out, "Folientext fehlt"
    assert "Geheime Sprechernotiz 555" in out, \
        "Sprechernotiz fehlt im Output (python-pptx-Extraktion geaendert?)"
    assert "Sprechernotizen" in r["note"], f"Note erwaehnt Sprechernotizen nicht: {r['note']}"


def build_structure_pptx(path):
    """Folien mit Bulletpoints (mehrstufig), Untertitel und Tabelle - deckt
    den python-pptx-Ersatz von markitdown ab (Struktur statt flachem Text)."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "Titelfolie"
    s1.placeholders[1].text = "Untertitel Gamma"
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Ergebnisse"
    body = s2.placeholders[1].text_frame
    body.text = "Oberpunkt"
    p = body.add_paragraph(); p.text = "Unterpunkt"; p.level = 1
    s3 = prs.slides.add_slide(prs.slide_layouts[5])
    s3.shapes.title.text = "Kennzahlen"
    tbl = s3.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(6), Inches(1)).table
    tbl.cell(0, 0).text = "Kennzahl"; tbl.cell(0, 1).text = "Wert"
    tbl.cell(1, 0).text = "Umsatz";   tbl.cell(1, 1).text = "1.2 Mio"
    prs.save(path)


def test_pptx_structure_preserved():
    """python-pptx-Ersatz: Folientrenner, Bullet-Hierarchie, Pipe-Tabelle,
    und Untertitel plain (KEINE erfundene Aufzaehlung)."""
    path = os.path.join(FIXTURES, "struktur.pptx")
    build_structure_pptx(path)
    r = convert_file(path, target_model="none")
    assert r["ok"], f"Konvertierung fehlgeschlagen: {r.get('error')}"
    out = r["output_text"]
    assert "## Folie 1: Titelfolie" in out, f"Folientrenner mit Titel fehlt: {out}"
    assert "- Oberpunkt" in out, "Bullet oberster Ebene fehlt"
    assert "  - Unterpunkt" in out, "Eingerueckte Bullet-Ebene fehlt (Hierarchie verloren)"
    assert "| --- | --- |" in out, "Pipe-Tabelle fehlt"
    assert "Untertitel Gamma" in out, "Untertitel fehlt"
    assert "- Untertitel Gamma" not in out, \
        "Untertitel faelschlich als Aufzaehlung ausgegeben"


def build_styleless_docx(path):
    """DOCX wie von Nicht-Word-Erzeugern (docx-js, pandoc-minimal): Absaetze
    ohne <w:pStyle> UND styles.xml ohne Default-Paragraph-Style. Erst diese
    Kombination macht paragraph.style zu None - fehlt nur das pStyle, faellt
    python-docx still auf den Default-Style ('Normal') zurueck."""
    import docx
    from docx.oxml.ns import qn
    doc = docx.Document()
    doc.add_heading("Protokoll Kickoff", level=1)
    doc.add_paragraph("Teilnehmer: Alle. Ort: Remote.")
    p = doc.add_paragraph("Absatz ohne jede Stil-Angabe.")
    pPr = p._p.find(qn("w:pPr"))
    if pPr is not None:
        ps = pPr.find(qn("w:pStyle"))
        if ps is not None:
            pPr.remove(ps)
    for st in doc.styles.element.findall(qn("w:style")):
        if (st.get(qn("w:type")) == "paragraph"
                and st.get(qn("w:default")) == "1"):
            st.set(qn("w:default"), "0")
    doc.save(path)


def test_docx_style_none_survives():
    """Block B: paragraph.style=None (Nicht-Word-DOCX) darf nicht crashen.
    Style-lose Absaetze werden als normaler Fliesstext behandelt."""
    path = os.path.join(FIXTURES, "ohne_style.docx")
    build_styleless_docx(path)
    r = convert_file(path)
    assert r["ok"], f"Konvertierung fehlgeschlagen: {r.get('error')}"
    out = r["output_text"]
    assert "# Protokoll Kickoff" in out, "Explizit gestylte Ueberschrift fehlt"
    assert "Absatz ohne jede Stil-Angabe." in out, "Style-loser Absatz fehlt"
    assert "# Absatz ohne jede Stil-Angabe." not in out, \
        "Style-loser Absatz faelschlich als Ueberschrift"


def test_csv_branch_keeps_companion_text():
    """Nebenpunkt: CSV-Zweig verwirft Begleittext nicht mehr."""
    path = os.path.join(FIXTURES, "tabellen.pdf")
    build_table_pdf_with_text(path)
    r = convert_file(path)
    assert r["ok"], f"Konvertierung fehlgeschlagen: {r.get('error')}"
    assert r["target_format"] == "csv", \
        f"Heuristik hat nicht CSV gewaehlt (got {r['target_format']}) - Testaufbau pruefen"
    out = r["output_text"]
    assert "R1C1" in out and "R10C5" in out, "Tabellendaten fehlen"
    assert "netto" in out, "Begleittext wurde im CSV-Zweig verworfen (alter Bug)"


def test_image_hint_in_text_pdf():
    """Nebenpunkt: eingebettete Bilder werden in der Note benannt."""
    path = os.path.join(FIXTURES, "mit_bild.pdf")
    build_image_in_text_pdf(path)
    r = convert_file(path)
    assert r["ok"], f"Konvertierung fehlgeschlagen: {r.get('error')}"
    assert "wichtiges Diagramm" in r["output_text"], "Text fehlt"
    assert "Bild" in r["note"], f"Note erwaehnt das nicht uebernommene Bild nicht: {r['note']}"


def test_regression_plain_text_pdf():
    """Regression: Standardfall (Text-PDF mit Kopf-/Fußzeilen) funktioniert weiter."""
    path = os.path.join(FIXTURES, "standard.pdf")
    build_plain_text_pdf(path)
    r = convert_file(path)
    assert r["ok"], f"Konvertierung fehlgeschlagen: {r.get('error')}"
    out = r["output_text"]
    for n in range(1, 4):
        assert f"Inhaltlicher Absatz Nummer {n}" in out, f"Inhalt von Seite {n} fehlt"
    assert "Firmenname AG - Interner Bericht" not in out, "Kopfzeile nicht entfernt"
    assert "Seite 2" not in out, "Fusszeile nicht entfernt"
    assert r["was_ocr"] is False, "Text-PDF darf kein OCR ausloesen"


def test_regression_full_image_pdf():
    """Regression: reine Bild-PDF laeuft weiter durch die komplette OCR."""
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader
    path = os.path.join(FIXTURES, "nur_bild.pdf")
    c = canvas.Canvas(path, pagesize=A4_PT)
    img = _make_text_image(["VOLLSTAENDIG GESCANNT: 2024"])
    c.drawImage(ImageReader(img), 0, 0, width=A4_PT[0], height=A4_PT[1])
    c.showPage()
    c.save()
    r = convert_file(path)
    assert r["ok"], f"Konvertierung fehlgeschlagen: {r.get('error')}"
    assert r["was_ocr"] is True, "Bild-PDF wurde nicht als OCR-Fall erkannt"
    assert "2024" in r["output_text"], "OCR-Text fehlt"
    assert "BILD-PDF" in r["note"], f"Note kennzeichnet Bild-PDF nicht: {r['note']}"
    # Block G1: reine Bild-PDF -> kein Pseudo-Prozentvergleich
    assert r["percent_saved"] is None, \
        f"OCR darf keinen Pseudo-Prozentwert liefern: {r['percent_saved']}"


# ---------------------------------------------------------------------------
# BLOCK G1: EINE TOKEN-METHODE PRO VERGLEICH
# ---------------------------------------------------------------------------

def test_g1_same_method_per_comparison():
    """Block G1: Vorher/Nachher desselben Vergleichs nutzen EINE Methode.
    Beleg pro Nicht-OCR-Dateityp-Pfad (txt, csv, xlsx, docx, pptx,
    Text-PDF): beide Seiten laufen durch count_tokens, token_method
    entspricht der global aktiven Methode."""
    expected = count_tokens("probe")["method"]

    txt_path = os.path.join(FIXTURES, "g1_probe.txt")
    with open(txt_path, "w", encoding="utf-8") as f:
        f.write("Zeile eins mit Inhalt.\nZeile zwei mit mehr Inhalt.")
    csv_path = os.path.join(FIXTURES, "g1_probe.csv")
    with open(csv_path, "w", encoding="utf-8") as f:
        f.write("Name,Wert\nAlpha,1\nBeta,2\n")
    xlsx_path = os.path.join(FIXTURES, "formeln.xlsx")
    build_formula_xlsx(xlsx_path)
    docx_path = os.path.join(FIXTURES, "reihenfolge.docx")
    build_order_docx(docx_path)
    pptx_path = os.path.join(FIXTURES, "notizen.pptx")
    build_notes_pptx(pptx_path)
    pdf_path = os.path.join(FIXTURES, "standard.pdf")
    build_plain_text_pdf(pdf_path)

    for path in (txt_path, csv_path, xlsx_path, docx_path, pptx_path, pdf_path):
        r = convert_file(path)
        assert r["ok"], f"{path}: {r.get('error')}"
        assert r["was_ocr"] is False, f"{path}: unerwartet OCR"
        assert r["token_method"] == expected, \
            (f"{os.path.basename(path)}: token_method {r['token_method']}, "
             f"erwartet {expected} (beide Seiten dieselbe Methode)")


def test_g1_estimate_fallback_both_sides():
    """Block G1: Ohne tiktoken fallen BEIDE Seiten gleichzeitig auf die
    Schaetzung zurueck (len/4) - nie gemischt."""
    import converter
    txt_path = os.path.join(FIXTURES, "g1_fallback.txt")
    content = "Zeile eins mit Inhalt.\nZeile zwei mit mehr Inhalt."
    with open(txt_path, "w", encoding="utf-8") as f:
        f.write(content)
    orig_enc, orig_tried = converter._TIKTOKEN_ENC, converter._TIKTOKEN_TRIED
    converter._TIKTOKEN_ENC, converter._TIKTOKEN_TRIED = None, True
    try:
        r = convert_file(txt_path, target_model="none")
        assert r["ok"]
        assert r["token_method"] == "estimate"
        assert r["tokens_before"] == max(1, round(len(content) / 4)), \
            "Vorher-Seite nutzt nicht die len/4-Schaetzung"
        assert r["tokens_after"] == max(1, round(len(r["output_text"]) / 4)), \
            "Nachher-Seite nutzt nicht die len/4-Schaetzung"
    finally:
        converter._TIKTOKEN_ENC, converter._TIKTOKEN_TRIED = orig_enc, orig_tried


# ---------------------------------------------------------------------------
# BLOCK G3: CSV-EXPORT OHNE \r\r-SEQUENZ
# ---------------------------------------------------------------------------

def test_g3_rows_to_csv_unified_newlines():
    """Block G3: _rows_to_csv liefert \\n wie alle uebrigen Exporte - kein
    \\r, das der Textmodus-Writer in app.py zu \\r\\r\\n verdoppeln wuerde."""
    from converter import _rows_to_csv
    text = _rows_to_csv([["a", "b"], ["c", "d"]])
    assert "\r" not in text, f"CR im CSV-String: {text!r}"
    assert text == "a,b\nc,d", f"Unerwartetes CSV: {text!r}"


def test_g3_csv_export_no_crcrlf_bytes():
    """Block G3: xlsx->CSV byte-genau ohne \\r\\r auf Platte. Der Schreibweg
    ist identisch zu app.py /convert (Textmodus, utf-8, Windows uebersetzt
    \\n -> \\r\\n)."""
    import openpyxl
    path = os.path.join(FIXTURES, "g3_export.xlsx")
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["Name", "Wert"])
    ws.append(["Alpha", 1])
    ws.append(["Beta", 2])
    wb.save(path)
    r = convert_file(path)
    assert r["ok"] and r["target_format"] == "csv"
    out_path = os.path.join(FIXTURES, "g3_export_out.csv")
    with open(out_path, "w", encoding="utf-8") as f:  # wie app.py:342
        f.write(r["output_text"])
    with open(out_path, "rb") as f:
        data = f.read()
    assert b"\r\r" not in data, f"\\r\\r-Sequenz in Datei-Bytes: {data[:80]!r}"
    assert b"Alpha,1" in data, "Inhalt fehlt"


# ---------------------------------------------------------------------------
# RUNNER
# ---------------------------------------------------------------------------

ALL_TESTS = [
    test_a_mixed_pdf_per_page_ocr,
    test_b_chapter_headings_protected,
    test_d_docx_original_order,
    test_e_xlsx_formula_fallback,
    test_pptx_notes_hint,
    test_pptx_structure_preserved,
    test_docx_style_none_survives,
    test_csv_branch_keeps_companion_text,
    test_image_hint_in_text_pdf,
    test_regression_plain_text_pdf,
    test_regression_full_image_pdf,
    test_g1_same_method_per_comparison,
    test_g1_estimate_fallback_both_sides,
    test_g3_rows_to_csv_unified_newlines,
    test_g3_csv_export_no_crcrlf_bytes,
]

if __name__ == "__main__":
    passed, failed = 0, 0
    for test in ALL_TESTS:
        name = test.__name__
        try:
            test()
            print(f"  PASS  {name}")
            passed += 1
        except Exception:
            print(f"  FAIL  {name}")
            traceback.print_exc()
            failed += 1
    print(f"\n{passed} bestanden, {failed} fehlgeschlagen.")
    sys.exit(1 if failed else 0)
